import os
import psutil
# file management
import glob
import os.path
from io import BytesIO
from datetime import date,datetime
from pytz import timezone as tz
import subprocess
import photoscript
import geopy,certifi,ssl
from geopy.geocoders import Nominatim

from deep_translator import GoogleTranslator

import re
import unicodedata
from unidecode import unidecode
from collections import defaultdict, Counter

#= Image Management ============================================
from PIL import Image, ImageDraw, ImageFont
from pillow_heif import register_heif_opener
import cv2
import numpy as np
#==== DOCX Section =============================================
from docx.shared import Mm,Pt,RGBColor 
#===============================================================
#-- _PPTX_
from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.enum.shapes import *
from pptx.enum.text import PP_ALIGN,MSO_AUTO_SIZE
from pptx.dml.color import *
from pptx.enum.dml import *
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.table import _Cell, _Row
from docx.oxml import parse_xml  ## for table backgroud
from docx.oxml.ns import nsdecls  ## for table bakcgroud
#import textwrap
import copy

# --------------------------------------------------------------
import requests # request img from web
import shutil # save img locally
# --------------------------------------------------------------
 
# --------------------------------------------------------------
icloud = "/Users/bernardconti/Library/Mobile Documents/com~apple~CloudDocs"
les_WIP = "/Users/bernardconti/LOCAL_TEMP/WIP/"
les_PDF_RV = "/Users/bernardconti/Library/Mobile Documents/com~apple~CloudDocs/Les éditions du 57/PDF_RectoVerso/"
icon_dir = icloud+'/MesProgrammes/python_global/icons'
export_path =  "/Users/bernardconti/LOCAL_TEMP/Photos"
save_path = "/Users/bernardconti/LOCAL_TEMP/Documents"
import_dir = "/Users/bernardconti/LOCAL_TEMP/ePhoto_Doc/Import/"
watermarkdir = "/Users/bernardconti/LOCAL_TEMP/watermark"
le_MH_Photos_Bio = icloud+'/MesProgrammes/MH_Photos_Bio/'
#--------------------------------------------------------------
retour_ligne = '\n'
Gray= "#EFEEEF"
GrayRow= "#EFEEEF"
Black ="#000000"
White = "#FFFFFF"
couleur_homme ="#ffcc99"
couleur_femme="#ccccff"
font_texte ="arial narrow"
couleur_cible = "#30DD30"
Gray1= "#b7b1b1"#F2F2F2
GraySide = "#c5c5c5"
Gray3 ="#D9D6D6"
Gray4 = "#E1DEDEEF"
couleur_chemin = "#EC5800"
couleur_titre ="#EB742F"
Green = "#B4E1C0"
Gray2= "#fde9d9"
Silver = "#A9A6A6" 
cell11 = "#F8CEB8"
cell12 = "#D9D9D9"
cell21 = "#FDE8DD"
cell22 = "#F2F2F2"
#Bleu = "#375e94"
Bleu = "#272170"
Darkblue = "#166082"
Gris_clair = "#EFEFEF"
ligne = "#375e94"
descendant = "#000000"
red = "#FF0000"
page_width = 190
pt_mm = 0.68
document_font = "Aptos (Corps)"
font_file = "/Users/bernardconti/Library/Fonts/Aptos.ttf"
font_file_bold = "/Users/bernardconti/Library/Fonts/Aptos-Bold.ttf"

EMU = 36000

#PPT
slide_width = 277
slide_height = 210
slide_margin_left = 5
slide_margin_right = 5
slide_margin_top = 10
slide_margin_bottom = 5
image_size = 28.2

slide_layout_box_0 = 0
slide_layout_box_0_1 = 1
slide_layout_ascendants_4 = 2
slide_layout_ascendants_tous = 3
slide_layout_couples_liste = 4
slide_layout_descendants_draw = 4
slide_layout_liste = 5
slide_layout_57 = 6
slide_layout_garde = 7
slide_layout_photoMH = 8
slide_layout_album = 9
slide_layout_entourage = 10
slide_layout_descendants_list = 11
slide_layout_section = 12
slide_layout_descendants_details = 13
slide_layout_table_image = 14
slide_layout_table_pleine = 15
#photos
box_photo_width = 50
box_photo_height = 50

list_events = ["Décoration","Distinction","Degree","Diplôme",
                "Military Service","Award","Honors","Title","Titre","Military Award",
                "Anoblissement","Nomination,Immigration","Association",
                "Language spoken","Illness","Comment","Marriage","Custom event","Nationalité",
                ]
#=============================================================================================================
# BOXES
#=============================================================================================================
def PPTX_add_box(la_page,x,y,w,h,*args,**kwargs):
#------------------------------------------------------------------------------------------------------------- 
    if w < 0 or h < 0 or not la_page :  
        print("PPTX_add_box",w,h)
        return False
#
# lecture des parametres
#
    la_couleur = False
    le_alignement = MSO_ANCHOR.MIDDLE
    la_margin_bottom = 0.5
    la_margin_top = 0.5
    la_margin_left = 0.5
    la_margin_right = 0.5
        
    for clef, valeur in kwargs.items():  
        if valeur: 
            if isinstance(clef,str): clef =clef.lower()
            if clef == "bcolor" : la_couleur = valeur
            if clef == "margin_bottom" : la_margin_bottom = valeur
            if clef == "margin_top" : la_margin_top = valeur
            if clef == "margin_left" : la_margin_left = valeur
            if clef == "margin_right" : la_margin_right = valeur

    isTransparent = False
    isWrap = True
    isAutosize = False
    isCadre = False

    for valeur in args:
        if valeur:
            if isinstance(valeur, str): valeur =valeur.lower()
            if valeur == "nowrap" : isWrap =  False    
            if valeur == "transparent" : isTransparent =  True  
            if valeur == "autosize" : isAutosize =  True    
            if valeur == "cadre" : isCadre =  True 
#
# ajout box
#
    txBox = la_page.shapes.add_textbox(Mm(x),Mm(y),Mm(w),Mm(h))

    if isCadre : 
        line = txBox.line
        line.color.rgb = RGBColor(200,200,200)
    elif la_couleur:
            txBox.fill.solid()
            txBox.fill.fore_color.rgb = hex_to_rgb(la_couleur)

    if isTransparent:
                            
        def SubElement(parent, tagname, **kwargs):
            element = OxmlElement(tagname)
            element.attrib.update(kwargs)
            parent.append(element)
            return element

        def _set_shape_transparency(shape, alpha):
        # Set the transparency (alpha) of a shape
            ts = shape.fill._xPr.solidFill
            sF = ts.get_or_change_to_srgbClr()
            sE = SubElement(sF, 'a:alpha', val=str(alpha))

        _set_shape_transparency(txBox,44000)

    tf = txBox.text_frame
    tf.word_wrap = isWrap
    tf.auto_size = isAutosize 
    if le_alignement : tf.vertical_anchor = le_alignement
    tf.margin_bottom = Mm(la_margin_bottom)
    tf.margin_top = Mm(la_margin_top)
    tf.margin_left = Mm(la_margin_left)
    tf.margin_right = Mm(la_margin_right)
#-------------------------------------------------------------------------------------------------------------
    return txBox
#=============================================================================================================
def PPTX_add_paragraph(text_Box,*args,**kwargs):
#------------------------------------------------------------------------------------------------------------- 
    align = PP_ALIGN.LEFT
    level = 0
    before_space = 0

    for valeur in args:
        if valeur:
            if isinstance(valeur, str): valeur =valeur.lower()
            if valeur == "center" : align =  PP_ALIGN.CENTER    
            if valeur == "right" : align =  PP_ALIGN.RIGHT           

    for clef, valeur in kwargs.items():  
        if valeur: 
            if isinstance(valeur,str): clef = clef.lower()
            if clef == "level" : level = int(valeur)
            if clef == "before_space" : before_space = int(valeur)
#-------------------------------------------------------------------------------------------------------------  
    tf = text_Box.text_frame
    if len(tf.paragraphs[0].text) == 0 : le_paragraph = tf.paragraphs[0] 
    else: le_paragraph = tf.add_paragraph()

    le_paragraph.level = level 
    le_paragraph.alignment = align
    le_paragraph.space_before = Pt(before_space)

    return le_paragraph
#=============================================================================================================
def PPTX_add_run(le_paragraph,le_texte,*args,**kwargs):
#------------------------------------------------------------------------------------------------------------- 
    run = None
    if le_texte:
#------------------------------------------------------------------------------------------------------------- 
        box_width = False
        isItalic =False
        isBold = False
        isUnderline = False
        font = False
        size = False
        color = False
        
        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "box_width" : box_width = valeur
                if clef == "size" : size = int(valeur)
                if clef == "font" : font = valeur
                if clef == "color" : color = valeur

        for valeur in args:
            if valeur:
                if isinstance(valeur, str): valeur =valeur.lower()
                if valeur == "bold" : isBold =  True    
                if valeur == "italic" : isItalic =  True  
                if valeur == "underline" : isUnderline =  True    
#-------------------------------------------------------------------------------------------------------------   
        if box_width:
            w_chr = (box_width-1)/len(le_texte)
            size = w_chr * 5
            if size > 12.5:size = 12.5
            if size < 8: size = 8
#------------------------------------------------------------------------------------------------------------- 
        run = le_paragraph.add_run()
        run.text =le_texte
        if font : run.font.name = font
        if size : run.font.size = Pt(size)
        if color : run.font.color.rgb = hex_to_rgb(color)        
        if isBold : run.font.bold = isBold
        if isItalic : run.font.italic = isItalic
        if isUnderline : run.font.underline = isUnderline

    return run
#=============================================================================================================
def PPTX_add_page(le_document,le_layout):
    boxes = []
    la_page = le_document.slides.add_slide(le_document.slide_layouts[le_layout])
    for shape in la_page.placeholders:
        #print('%d %s' % (shape.placeholder_format.idx, shape.name))
        boxes.append(shape)
    b = PPTX_add_box(la_page,slide_width-10,slide_height-5,10,5)
    n_slide = le_document.slides.index(la_page)
    if n_slide > 1 : PPTX_add_run(PPTX_add_paragraph(b,"right"),str(n_slide),"italic","bold",size=10)
    return la_page,boxes
#=============================================================================================================
def hex_to_rgb(value):
#-------------------------------------------------------------------------------------------------------------
    RGB = RGBColor(0, 0, 0)
    try:
        value = value.lstrip('#')
        lv = len(value)
        (rouge, vert, bleu) = tuple(int(value[i:i + lv // 3], 16) for i in range(0, lv, lv // 3))
        RGB = RGBColor(rouge, vert, bleu)
    except Exception as error:
        print (value,tuple(int(value[i:i + lv // 3], 16) for i in range(0, lv, lv // 3)))
#-------------------------------------------------------------------------------------------------------------
    return RGB
#=============================================================================================================
def watermark(image_path,le_texte): 
# --------------------------------------------------------------   
#Import required Image library
    #Create an Image Object from an Image
    if le_texte : 
        im = Image.open(image_path)
        width, height = im.size

        final_im = Image.new('RGB', (width, int(height*1.11)))
        watermark_im = Image.new('RGB', (width, int(height*0.11)),color = (256, 256, 256))
        draw = ImageDraw.Draw(watermark_im)
        fsize = height * 0.1
        font = ImageFont.truetype("Monaco.ttf",fsize)

        # calculate the x,y coordinates of the text
        x = width/10
        y =0

        position = (x, y)
        left, top, right, bottom = draw.textbbox(position, le_texte, font=font)
        #draw.rectangle((left-10, top-10, right+10, bottom+10), fill="white")
        #draw.rectangle((left, top, right, bottom), fill="white")
        draw.text(position, le_texte, font=font, fill="black")
        #draw.text(position, le_texte, font=font, fill="white")

        #im.show()
        w_file_name = watermarkdir+"/w_"+image_path.split("/")[-1]

        #Save watermarked image
        final_im.paste(im, (0,0))
        final_im.paste(watermark_im, (0,height))

        final_im.save(w_file_name)
    else: w_file_name = image_path

    return w_file_name
#=============================================================================================================

#====PDF =====================================================================================================
def PDF_add_page(le_pdf,image_filename,n_page):
#=============================================================================================================
    n_page = n_page + 1
    le_pdf.add_page()
    
    if n_page % 2 == 0: 
        xi = 0
        xt = 282
        xc = 287
    else: 
        xi = 20
        xt = 5
        xc = 10
    
    if n_page > 1:
        le_pdf.set_font('helvetica', size=12)
        le_pdf.text(x= xt ,y=205,text=f'# {n_page}')
    else:        
        le_pdf.set_font('helvetica', size=10)
        le_pdf.text(x= xt ,y=200,text=f'Édition')
        le_pdf.text(x= xt ,y=205,text=f'du 57')

    le_pdf.circle(xc,55,2.5,"F")
    le_pdf.circle(xc,155,2.5,"F")

    if image_filename :  le_pdf.image(image_filename,  x = xi , y = 0 ,w= 277,h=210)

    return n_page
#=============================================================================================================
def excute_cmd(command): 
# --------------------------------------------------------------   
    command_list = command.split(" ")
    #print(command_list)
    result = subprocess.run(command_list, capture_output=True, text=True)
    if result.stdout : print ("resultat = " + result.stdout)
    if result.stderr : print("erreur = " + result.stderr)
# --------------------------------------------------------------
    return(result.stdout,result.stderr)
#=============================================================================================================

#=============================================================================================================
def PPTX_expand_picture(pic,x_bandeau) :
        if pic:
            pvh = pic.height
            r = pic.width / pvh

            pic.width = (slide_width - x_bandeau ) * EMU - pic.left 
            pic.height = int(pic.width / r)

            pic.crop_bottom = -(1-pvh/pic.height)
            pic.crop_top = (1 - pvh/pic.height)
        return
#=============================================================================================================

#===========================================================================================
# DEBUT
#===========================================================================================
print("Execution de ",__file__)
dir  = "/Users/bernardconti/Downloads/"

#============================================================================================   
#   Générationdu du type d'édition
#============================================================================================
#--- Initialisation avec les données de la cible
#----------------------------------------------------------------------------------------------
#-- MENU ---------------------------------------------------------------------------------------
from PyQt6.QtWidgets import *
from PyQt6.QtGui import *
from PyQt6.QtCore import *

custom_font = QFont("Arial",14)
line_height = 30
line_spacing = 10
cell_spacing = 5
n_line = 7
bio_width = [400,120,120,120,120,120]
bio_layout_width = 0
for number in bio_width:
    bio_layout_width += number
bio_layout_width += (len(bio_width)- 1)*cell_spacing
bio_layout_height = n_line * line_height + (n_line-1)*line_spacing

cell_width = 250

#========================================================================================    
for proc in psutil.process_iter():
    if proc.name() == "Microsoft PowerPoint": proc.kill()
#========================================================================================
# recupere tous les noms
#========================================================================================

# initiate list des directore pour impression PDF R&V
list_dirs = glob.glob(les_WIP+"*/", recursive = False)
if not list_dirs : exit()
#========================================================================================   
#  MAIN WINDOW   
#========================================================================================
class LA_WINDOW(QMainWindow):
#========================================================================================
    def __init__(self):
        #super().__init__(*args, **kwargs)
        super().__init__()

        #try:
        #    la_db = "/Users/bernardconti/Pictures/Photos Library.photoslibrary/database/Photos.sqlite"
        #    import sqlite3
        #    connection = sqlite3.connect(la_db, timeout=10)
        #    cur = connection.cursor()
        #    cur.execute("SELECT name FROM sqlite_master WHERE type='table' ORDER BY name;")
        #    #cur.execute('''create table item
        #    #(id integer primary key, itemno text unique,
        #    #        scancode text, descr text, price real)''')#

        #connection.commit()
        #    cur.close()
        #except Exception as error:
        #    print(error)

        import osxphotos 
        try:
            self.photosdb = osxphotos.PhotosDB()
            #self.photosdb = osxphotos.PhotosDB(dbfile = "/Users/bernardconti/Pictures/PhotosLibrary_BC.photoslibrary")
            isOsxPhoto = True
            #print("OK Osxphotos")
        #"/Users/bernardconti/Pictures/Photos Library.photoslibrary"
        except Exception as error:
            print("Base de donnée ",error)
            isOsxPhoto = False

        self.setupUi()
#----------------------------------------------------------------------------------------
    def central_accueil(self):
        self.menu_clear_layout(self.central_layout)
        for item in self.infos:
            info = QLabel(item)
            info.setFixedSize(bio_layout_width, line_height)
            info.setStyleSheet("background-color: lightgray;color: black;"
                                    "qproperty-alignment: AlignLeft;"
                                    "border-radius: 10px;"
                                    "qproperty-wordWrap: true;"
                                    "padding: 6px;"
                                    )
            info.setFont(custom_font)
            self.central_layout.addRow(info) 
        self.menu_add_infos_row("2025, Created by Bernard CONTI, Dourdan, France")
#===========================================================================================================
    def setupUi(self):

        self.setObjectName("Menu_principal")
        self.resize(400, 400)

# Create menu bar, toolbar and statusbar objects
        menubar = self.menuBar()
        menubar.setFont(custom_font)
        
        self.centralwidget = QWidget(self)
        self.setCentralWidget(self.centralwidget)
        self.statusbar = QStatusBar(self)

# Create actions
        finAction = self.menu_add_action(icon = "exit.png",text = "Au revoir et Merci", cmd = self.PPTX_fin)
        #updatePhotos_album_Action = self.menu_add_action(text = "Update Photos Album", cmd = self.central_updatePhotos_album)
        importPhotos_Action = self.menu_add_action(text = "Import Photos", cmd = self.central_import_photos)
        updatePhotos_selection_Action = self.menu_add_action(text = "Update Photos Selection", cmd = self.central_updatePhotos_selection)
        negatifsAction = self.menu_add_action(text = "Négatifs", cmd = self.central_negatifs)

        pdfAction = self.menu_add_action(text = "PDF Recto Verso", cmd = self.central_pdf)
        livresAction = self.menu_add_action(text = "Album", cmd = self.central_livres)
        BienvenueAction = self.menu_add_action(icon = "home.png", text = "Home Sweet Home",cmd = self.central_accueil)
# Creation des mennus     
        helpMenu = menubar.addMenu(QIcon(icon_dir+"/tools.png"),"&Outils")
        helpMenu.setStyleSheet("background-color: "+ couleur_homme +";color: black;"
        "font-name=Arial;"
        "font-size: 16px;"
        )
# Add actions to menus
        #helpMenu.addAction(updatePhotos_album_Action)
        self.last_date = None
        self.last_heure = None
        helpMenu.addAction(importPhotos_Action)
        helpMenu.addAction(updatePhotos_selection_Action)
        helpMenu.addAction(negatifsAction)
        helpMenu.addAction(pdfAction)

# add menus to menubar
        menubar.addMenu(helpMenu)

# add  toolbar
        toolbar = QToolBar('Main ToolBar', self)
        toolbar.setIconSize(QSize(24, 24))
        toolbar.setStyleSheet("background-color: "+ couleur_femme +";color: black;"
                    "font-name=Arial;"
                    "font-size: 16px;"
                    )

        toolbar.addAction(BienvenueAction)
        toolbar.addAction(livresAction)
        toolbar.addActions(menubar.actions())
        toolbar.addAction(finAction)
        
# add toolbar and statusbar to main window
        self.addToolBar(toolbar)
        self.setStatusBar(self.statusbar)

# central_layout, self
        self.central_layout = QFormLayout(self.centralwidget) 
        self.central_layout.setVerticalSpacing(line_spacing)
        self.central_layout.setHorizontalSpacing(cell_spacing) 
#init central on welcome
        
        self.infos = ["Bienvenue dans le MIFA Studio"]
        self.central_accueil()

#======================================================================================================================          
    def PPTX_save(self):
        la_page = self.le_document.slides.add_slide(self.le_document.slide_layouts[slide_layout_57])
        OUT_fichier  = f'{save_path}/{self.OUTFICHIER.text()}'
        self.le_document.save(OUT_fichier)

        for file in os.listdir(watermarkdir):
            try: os.remove(watermarkdir+"/"+file)
            except Exception as error: print(error)

        self.infos.append("Livres > "+ OUT_fichier +  " enregistré")
        self.central_accueil()

#==============================================================================================================            
    def PPTX_fin(self):
        print("Au revoir et merci")
        self.force_close = True
        self.close()
#==============================================================================================================

#==============================================================================================================
# DEF_menu_add_...
#============================================================================================================== 
    def list_albums(self,le_layout,cmd):
#==============================================================================================================
        self.menu_clear_layout(le_layout)

        self.next = self.menu_add_pushbutton('Suite',cmd,n_cell = 2)


        self.tree = QTreeWidget()
        self.tree.setHeaderHidden(True)
        self.tree.setFixedSize(2*cell_width, 10* line_height)

        selmodel = self.tree.selectionModel()
        selmodel.selectionChanged.connect(self.click_action)

        le_layout.addRow(self.tree)
        le_layout.addRow(self.next)

        self.list_selection_albums = []
        self.les_folders = []
        for f in self.photosdb.folder_info:
            if f.title == "LES LIVRES_CONTI": 
                for fs in f.subfolders:

                    isP = True
                    for le_album in fs.album_info:
                        if le_album.title[0] != "P":
                            isP = False
                            break

                    if not isP:
                        self.les_folders.append([fs.title,fs])
                        level0 = QTreeWidgetItem(self.tree, [fs.title])
                        for le_album in fs.album_info:
                            if le_album.title[0] != "P": level1 = QTreeWidgetItem(level0, [le_album.title])
                    
#===========================================================================================
    def click_action(self, selected, deselected):

        for index in selected.indexes():
            item = self.tree.itemFromIndex(index)
            if self.tree.itemFromIndex(index.parent()) : 
                if item.text(0) not in self.list_selection_albums : 
                    self.list_selection_albums.append(item.text(0))
                    self.menu_add_infos_row(f'{self.tree.itemFromIndex(index.parent()).text(0)} > {item.text(0)} ')
            else:
                for f in self.les_folders:
                    if item.text(0) == f[0]:
                        for le_album in f[1].album_info:
                            if le_album.title not in self.list_selection_albums : 
                                if le_album.title[0] != "P": self.list_selection_albums.append(le_album.title)
                        self.menu_add_infos_row(f'Dossier : {item.text(0)} ')

#================================================================================================
    def livres_tab_album(self):
        self.list_albums(self.livres_album_layout,self.livres_tab_album_exec)
#================================================================================================
    def livres_tab_album_exec(self):

        if not self.list_selection_albums : return

        self.menu_clear_layout(self.livres_album_layout)

        self.livres_album_layout.addRow(self.menu_add_label("Albums",n_cell = 2,bcolor = "green",color = "white"))
        for a in self.list_selection_albums:
            self.livres_album_layout.addRow(self.menu_add_label(a,n_cell = 2))
# section regroupement      
        self.livres_album_layout.addRow(self.menu_add_label("Configuration du séquencement :",n_cell = 2,bcolor = "green",color = "white"))
# type de tri
        self.type_tri  = self.menu_add_input_text_combobox(["Date","Titre","Fichier"],n_cell = 1)
        self.livres_album_layout.addRow(self.menu_add_label("Type de tri des photos"),self.type_tri)
# Changement de slide
        self.CesureDate = self.menu_add_checkbox("Oui",n_cell = 1)
        self.livres_album_layout.addRow(self.menu_add_label("Césure Date"),self.CesureDate)
# Changement de slide
        self.CesureVille = self.menu_add_checkbox("Oui",n_cell = 1)
        self.livres_album_layout.addRow(self.menu_add_label("Césure Ville"),self.CesureVille)
# Changement de slide
        self.CesureComment = self.menu_add_checkbox("Oui",n_cell = 1)
        self.livres_album_layout.addRow(self.menu_add_label("Césure Commentaire"),self.CesureComment)
# section Marge Haute      
        self.livres_album_layout.addRow(self.menu_add_label("Configuration de la marge haute :",n_cell = 2,bcolor = "green",color = "white"))
# Marge haute
        self.marge_haute = self.menu_add_input_text_combobox(["Sans","Petite","Moyenne","Grande"],n_cell = 1)
        self.livres_album_layout.addRow(self.menu_add_label("Marge Haute"),self.marge_haute)

# section watermark      
        self.livres_album_layout.addRow(self.menu_add_label("Configuration du Watermark :",n_cell = 2,bcolor = "green",color = "white"))
# Info photos       
        self.avec_personnes = self.menu_add_input_text_combobox(["Sans","Avec"],n_cell = 1)
        self.livres_album_layout.addRow(self.menu_add_label("Les personnes"),self.avec_personnes)
# Info Dates     
        self.avec_date = self.menu_add_input_text_combobox(["Sans","Avec"],n_cell = 1)
        self.livres_album_layout.addRow(self.menu_add_label("Les dates"),self.avec_date)
# Info Ville  
        self.avec_ville = self.menu_add_input_text_combobox(["Sans","Avec"],n_cell = 1)
        self.livres_album_layout.addRow(self.menu_add_label("Les villes"),self.avec_ville)
# Info Fichier  
        self.avec_fichier = self.menu_add_input_text_combobox(["Sans","Avec"],n_cell = 1)
        self.livres_album_layout.addRow(self.menu_add_label("Les fichier"),self.avec_fichier)
# Nom de l'album photos à publier
#----------------------------------------------------------------------------------------------------            
# Action de Publication
        self.livres_album_layout.addRow(self.menu_add_pushbutton("Publish Albums",self.PPTX_publication_album))      
#==============================================================================================================
    def livres_tab_document(self):
        self.menu_clear_layout(self.livres_document_layout)          
#--------------------------------------------------------------------------------------------------------------
# recupere tous les documents
        temp_documents = []
        for item in os.listdir(f'{save_path}/'):
            if item.split(".")[-1].lower() == "pptx" and "$" not in item:
                temp_documents.append(item)
        temp_documents = sorted(temp_documents, key=lambda col: (col[0]) )
        list_documents = ["Nouveau"] + temp_documents
#--------------------------------------------------------------------------------------------------------------
# recupere tous les models
        self.PPTX_label = self.menu_add_label("Choisir un fichier existant ou Initialiser un Nouveau PPTX",n_cell = 3)
        self.PPTX_file = self.menu_add_input_text_combobox(list_documents)  
        self.PPTX_action = self.menu_add_pushbutton("Suite ...",self.PPTX_Document)
        self.livres_document_layout.addRow(self.PPTX_label)
        self.livres_document_layout.addRow(self.PPTX_file,self.PPTX_action)
#===========================================================================================================
    def central_livres(self):
        self.menu_clear_layout(self.central_layout)

        self.livres_document,self.livres_document_layout = self.menu_add_layout("form")

        self.livres_tab = QTabWidget() 
        self.livres_tab.addTab(self.livres_document, 'Document')
        self.central_layout.addRow(self.livres_tab)

        self.livres_tab_document()
#===========================================================================================================
    def central_pdf(self):
        self.menu_clear_layout(self.central_layout)
        self.PDF_leWIP = self.menu_add_input_text_combobox(list_dirs,n_cell = 2)
        self.pdf_action = self.menu_add_pushbutton("Générer le PDF Recto Verso",self.PDF_generate)
        self.central_layout.addRow(self.PDF_leWIP,self.pdf_action)
#===========================================================================================================           
    def central_negatifs(self):
        self.menu_clear_layout(self.central_layout)       
        
        save_path = "/Users/bernardconti/LOCAL_TEMP/Négatifs/"
        target_album="NÉGATIFS_CONVERTIS"

        for idx,photo in enumerate(self.photosdb.photos(albums=["NÉGATIFS_TEMP"])):

            image = cv2.imread(photo.path_derivatives[0])
            gray_image = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
            negative_image = 255 - gray_image 
            ct=datetime.now()
            savefilename = f'{save_path}NEGATIF_{ct.year}-{ct.month}-{ct.day}_{ct.hour}:{ct.minute}:{ct.second}_{idx}.jpg'       
            self.central_layout.addRow(self.menu_add_label(f'NEGATIF_{ct.year}-{ct.month}-{ct.day}_{ct.hour}:{ct.minute}:{ct.second}_{idx}.jpg',n_cell = 2 ))
            cv2.imwrite(savefilename, negative_image) 
            command_list = ["osxphotos","import",savefilename,"--album",target_album,"--skip-dups"]
            result = subprocess.run(command_list, capture_output=False, text=True)
            os.remove(savefilename)
#===========================================================================================================
    def central_updatePhotos_selection(self):
        self.menu_clear_layout(self.central_layout) 
        self.central_layout.addRow(self.menu_add_pushbutton("Sélection des photos",self.central_updatePhotos_selection_exec, n_cell = 2, bcolor = "green", color = "white"))
#===========================================================================================================
    def central_updatePhotos_selection_exec(self):

        photoslib = photoscript.PhotosLibrary()
        if photoslib.selection:
            self.photos_selection = photoslib.selection
            for la_photo in self.photos_selection:
                self.central_layout.addRow(self.menu_add_label(f'{la_photo.filename}',n_cell = 2))

        self.la_date_photos = self.menu_add_date(default = self.last_date,ncol = 2)
        self.la_heure_photos = self.menu_add_heure(default = self.last_heure,ncol = 2)
        self.le_écart_photos =self.menu_add_input_text_combobox(["10 minutes", "30 minutes", "1 heure"],n_cell = 1)
        self.le_commentaire_photos =self.menu_add_input_text_simple("Sans",n_cell = 1)
        self.le_lieu_photos = self.menu_add_input_text_simple("Sans",n_cell = 1)

        self.central_layout.addRow(self.menu_add_label("Jour"),self.la_date_photos)
        self.central_layout.addRow(self.menu_add_label("Heure"),self.la_heure_photos)
        self.central_layout.addRow(self.menu_add_label("Écart"),self.le_écart_photos)
        self.central_layout.addRow(self.menu_add_label("Commentaire)"),self.le_commentaire_photos)
        self.central_layout.addRow(self.menu_add_label("Lieu"),self.le_lieu_photos)

        self.central_layout.addRow(self.menu_add_pushbutton("Metre à jour ...",self.PHOTOS_UpdatePhotos_selection, n_cell = 1))
#===========================================================================================================
    def central_import_photos(self):
        self.menu_clear_layout(self.central_layout)       
        self.list_albums(self.central_layout,self.central_import_photos_exec) 
#===========================================================================================================               
    def central_import_photos_exec(self):
        self.menu_clear_layout(self.central_layout)
#----------------------------------------------------------------------------------------------------
        self.la_date_photos = self.menu_add_date(default = self.last_date,ncol = 2)
        self.la_heure_photos = self.menu_add_heure(default = self.last_heure,ncol = 2)
        self.le_écart_photos =self.menu_add_input_text_combobox(["10 minutes", "30 minutes", "1 heure"],n_cell = 1)
        self.le_commentaire_photos =self.menu_add_input_text_simple("Sans",n_cell = 1)
        self.le_lieu_photos = self.menu_add_input_text_simple("Sans",n_cell = 1)

        self.central_layout.addRow(self.menu_add_label("Jour"),self.la_date_photos)
        self.central_layout.addRow(self.menu_add_label("Heure"),self.la_heure_photos)
        self.central_layout.addRow(self.menu_add_label("Écart"),self.le_écart_photos)
        self.central_layout.addRow(self.menu_add_label("Commentaire)"),self.le_commentaire_photos)
        self.central_layout.addRow(self.menu_add_label("Lieu"),self.le_lieu_photos)

        self.central_layout.addRow(self.menu_add_pushbutton("Metre à jour ...",self.PHOTOS_Import, n_cell = 1))
#===========================================================================================================   
    def menu_add_layout(self,type):
#--------------------------------------------------------------------------------------------------------------
        if type == "grid": item_layout = QGridLayout()
        else: item_layout = QFormLayout()

        item_layout.setVerticalSpacing(line_spacing)
        item_layout.setHorizontalSpacing(cell_spacing)  
        item = QWidget(self)
        item.setLayout(item_layout)
        return item,item_layout
#============================================================================================================== 
    def menu_clear_layout(self,le_layout):
#--------------------------------------------------------------------------------------------------------------
        while le_layout.count() > 0:
            item = le_layout.itemAt(0)
            widget = item.widget()
            if widget is None:
                le_layout.removeItem(item)
            else:
                #print(widget)
                widget.deleteLater()
                le_layout.removeWidget(widget)
        #print(le_layout.count())
        return
#==============================================================================================================         
    def menu_add_input_text_simple(self,default_input_text,**kwargs):
#--------------------------------------------------------------------------------------------------------------
        if not default_input_text : return
        la_width = None
        n_cell = 2
        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "width" : la_width = valeur
                if clef == "n_cell" : n_cell = valeur
        if not la_width : la_width = n_cell * cell_width
#-------------------------------------------------------------------------------------------------------------- 
        item = QLineEdit(default_input_text)     
        item.setStyleSheet("background-color: "+couleur_femme+";color: black")
        item.setAlignment(Qt.AlignmentFlag.AlignHCenter)
        item.setFont(custom_font) 
        item.setFixedSize(la_width,line_height) 
        return item
#==============================================================================================================
    def menu_add_input_text_completor(self,default_input_list,**kwargs):
#--------------------------------------------------------------------------------------------------------------
        if not default_input_list : return
        la_width = None
        le_alignment = "center"
        n_cell = 2
        popup_action = None
        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "width" : la_width = valeur
                if clef == "n_cell" : n_cell = valeur
                if clef == "popup"  : popup_action = valeur
                if clef == "alignment": le_alignment = valeur
        if not la_width : la_width = n_cell * cell_width

#--------------------------------------------------------------------------------------------------------------    
        le_completor = QCompleter(default_input_list)
        le_completor.setCaseSensitivity(Qt.CaseSensitivity.CaseInsensitive)
        if popup_action : le_completor.popup().clicked.connect(popup_action)
        le_completor.popup().setStyleSheet(
            "QListView {background-color: "+couleur_femme+";color: black;"
                        "selection-background-color: "+couleur_homme+";}"
            )
        item = QLineEdit(self)            
        item.setCompleter(le_completor)
        item.setStyleSheet("background-color: "+couleur_femme+";color: black")
        
        if le_alignment == "left":
            item.setAlignment(Qt.AlignmentFlag.AlignLeft)
        elif le_alignment == "right":
            item.setAlignment(Qt.AlignmentFlag.AlignRight)
        else:
            item.setAlignment(Qt.AlignmentFlag.AlignHCenter)
        item.setFont(custom_font) 
        item.setFixedSize(la_width,line_height) 
        
        return item     
#==============================================================================================================
    def menu_add_input_text_combobox(self,default_input_list,**kwargs):
#--------------------------------------------------------------------------------------------------------------
        if not default_input_list : return
        la_width = None
        n_cell = 2
        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "width" : la_width = valeur
                if clef == "n_cell" : n_cell = valeur
        if not la_width : la_width = n_cell * cell_width
#--------------------------------------------------------------------------------------------------------------
        item = QComboBox()
        item.addItems(default_input_list)      
        item.setStyleSheet("background-color: "+couleur_homme+";color: black")
        item.setFont(custom_font) 
        item.setFixedSize(la_width,line_height) 

        return item
#==============================================================================================================
    def menu_add_infos_row(self,le_texte):
#--------------------------------------------------------------------------------------------------------------
        self.statusbar.showMessage(le_texte)
        self.statusbar.setStyleSheet("background-color: "+couleur_homme+";color: black")
#==============================================================================================================
    def menu_add_label(self,le_label,**kwargs):
#--------------------------------------------------------------------------------------------------------------
        la_width = None
        bcolor = "lightgray"
        color = "black"
        n_cell = 1
        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "width" : la_width = valeur
                if clef == "n_cell" : n_cell = valeur
                if clef == "bcolor": bcolor = valeur
                if clef == "color": color = valeur

        if not la_width : la_width = n_cell * cell_width
#--------------------------------------------------------------------------------------------------------------                   
        item = QLabel(le_label)
        item.setStyleSheet("background-color: "+ bcolor +";color: "+color+";"
                                "qproperty-alignment: AlignLeft;"
                                "border-radius: 10px;"
                                "qproperty-wordWrap: true;"
                                "padding: 6px;"
                                )
        item.setFont(custom_font) 
        #item_label.setAlignment(Qt.AlignmentFlag.AlignHCenter)
        item.setFixedSize(la_width,line_height)

        return item
#==============================================================================================================
    def menu_add_date(self,**kwargs):
        la_width = None
        n_cell = 1
        format_osxphotos = "yyyy-MM-dd"
        format_edit = "dd-MM-yyyy"
        default = False

        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "width" : la_width = valeur
                if clef == "n_cell" : n_cell = valeur
                #if clef == "format":  format = valeur
                if clef == "default" : default = valeur
        if not la_width : la_width = n_cell * cell_width
#--------------------------------------------------------------------------------------------------------------
        item =QDateEdit()
        item.setDisplayFormat(format_edit)
        if default :
            item.setDate(QDate.fromString(default,format_edit))
        else :
            item.setDate(QDate.currentDate())
            
        item.setFont(custom_font)
        item.setFixedSize(la_width, line_height)
        item.setAlignment(Qt.AlignmentFlag.AlignHCenter)

        return item
#==============================================================================================================
    def menu_add_heure(self,**kwargs):
        la_width = None
        n_cell = 1
        default = False

        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "width" : la_width = valeur
                if clef == "n_cell" : n_cell = valeur
                if clef == "default" : default = valeur
        if not la_width : la_width = n_cell * cell_width
#--------------------------------------------------------------------------------------------------------------
        item =QTimeEdit()
        item.setDisplayFormat('HH:mm')
        if default :
            item.setTime(QTime.fromString(default,'HH:mm:ss'))
        else :
            #item.setTime(QTime.currentTime())
            item.setTime(QTime(10,00,00))
            
        item.setFont(custom_font)
        item.setFixedSize(la_width, line_height)
        item.setAlignment(Qt.AlignmentFlag.AlignHCenter)

        return item
#==============================================================================================================
    def menu_add_pushbutton(self,le_texte,cmd,**kwargs):
#--------------------------------------------------------------------------------------------------------------
        la_width = None
        n_cell = 1
        bcolor = "green"
        color = "white"

        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "width" : la_width = valeur
                if clef == "n_cell" : n_cell = valeur
                if clef == "color" : color = valeur
                if clef == "bcolor" : bcolor = valeur
        if not la_width : la_width = n_cell * cell_width
#--------------------------------------------------------------------------------------------------------------
        item = QPushButton(le_texte)
        item.setStyleSheet(
                "QPushButton{background-color : "+bcolor+";color: "+color+"}"
                "QPushButton::pressed{background-color : red ;color: white}"
                        )
        item.setFont(custom_font)
        item.clicked.connect(cmd)
        item.setFixedSize(la_width, line_height)

        return item
#==============================================================================================================
    def menu_add_checkbox(self,le_texte,**kwargs):
#--------------------------------------------------------------------------------------------------------------
        la_width = None
        n_cell = 1
        color_check = couleur_homme
        color = couleur_femme

        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "width" : la_width = valeur
                if clef == "n_cell" : n_cell = valeur
                if clef == "color" : color = valeur
                if clef == "bcolor" : bcolor = valeur
        if not la_width : la_width = n_cell * cell_width
#--------------------------------------------------------------------------------------------------------------
        item = QCheckBox(le_texte)
        item.setStyleSheet(
            "QCheckBox{background-color : "+color+";color:black;}"
            "QCheckBox::checked{background-color : "+color_check+";color:black;}" )

        item.setFont(custom_font)
        item.setFixedSize(la_width, line_height)

        return item
#==============================================================================================================   
    def menu_add_action(self,**kwargs):
#--------------------------------------------------------------------------------------------------------------
        tip = None
        cmd = None
        icon = None
        text = None

        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "tip": tip = valeur
                if clef == "cmd": cmd = valeur
                if clef == "icon": icon = valeur
                if clef == "text": text = valeur
#--------------------------------------------------------------------------------------------------------------

        if icon and text :       item = QAction(QIcon(icon_dir+"/"+icon), text, self)
        elif not icon and text : item = QAction(text, self)
        elif icon and not text : item = QAction(QIcon(icon_dir+"/"+icon), "-",self)
        else: return None

        if cmd : item.triggered.connect(cmd)
        if tip : 
            #item.setToolTip(tip)
            item.setStatusTip(tip)
        item.setCheckable(True)
    
        return item                        
#==============================================================================================================
    def les_sub_folders(self,f_cur):
        temp_list = []

        for le_album in f_cur.album_info:

            if le_album.title[0] != "P": temp_list.append(f' > {le_album.title}')

        if temp_list   :   
            self.list_selection_albums.append(f_cur.title)
            self.list_folder_albums.append([f_cur.title,temp_list])
            self.list_selection_albums = self.list_selection_albums + temp_list

        for sf in f_cur.subfolders:
            self.les_sub_folders(sf)
        return 
#==============================================================================================================           
    def PPTX_reset(self):

        for item in self.personnes_options:
            for iitem in item:
                if iitem:
                    try: iitem.deleteLater() 
                    except : next 
        try:
            self.publication.deleteLater()
            self.mode_couples.deleteLater()
            self.mode_fratrie.deleteLater()
        except : next 
#==============================================================================================================                       
    def PPTX_Document(self):
#-------------------------------------------------------------------------------------------------------------- 
        le_fichier = self.PPTX_file.currentText()

        if le_fichier == "Nouveau":       
            ct=datetime.now()
            self.OUTFICHIER = f'LIVRE_{ct.year:04d}{ct.month:02d}{ct.day:02d}{ct.hour:02d}{ct.minute:02d}{ct.second:02d}.pptx'
            # recupere tous les models

            list_models = []
            for item in os.listdir(f'{icloud}/MesProgrammes/Mes_Models/'):
                if item.split(".")[-1].lower() == "pptx" and item.startswith('MODEL_') :
                    list_models.append(item.replace(".pptx",""))
            list_models = sorted(list_models, key=lambda col: (col[0]) )
        
            if not list_models : 
                print("Pas de MODEL_")
                exit()
            
            self.livres_document_layout.addRow(self.menu_add_label("Model actif pour la création de Livre"))
            self.model_file = self.menu_add_input_text_combobox(list_models)
            self.livres_document_layout.addRow(self.model_file,self.menu_add_pushbutton("Suite ...",self.PPTX_add_new_document))
                                            
        else:

            IN_fichier  = f'{save_path}/{le_fichier}'
            self.OUTFICHIER = le_fichier
            print(IN_fichier)
            try :
                self.le_document = Presentation(IN_fichier)
                os.remove(IN_fichier)
            except Exception as error: 
                print(error)
                exit()

            self.menu_clear_layout(self.livres_document_layout)
            self.PPTX_label2 = self.menu_add_label("Vous pouvez modifier le nom du livre",n_cell = 3)
            self.livres_document_layout.addRow(self.PPTX_label2)
            self.OUTFICHIER = self.menu_add_input_text_simple(self.OUTFICHIER)
            self.PPTX_action = self.menu_add_pushbutton("Enregistrer le livre",self.PPTX_save)
            self.livres_document_layout.addRow(self.OUTFICHIER,self.PPTX_action)

            self.livres_bio,self.livres_bio_layout = self.menu_add_layout("grid")
            self.livres_tab.addTab(self.livres_bio,"Biographie")
            self.livres_tab_bio()

            self.livres_album, self.livres_album_layout = self.menu_add_layout("from")
            self.livres_tab.addTab(self.livres_album,"Album")
            self.livres_tab_album()

            self.livres_tab.setCurrentIndex(1)
    #==============================================================================================================   
    def PPTX_add_new_document(self):
        
            self.le_document = Presentation(f'{icloud}/MesProgrammes/Mes_Models/{self.model_file.currentText()}.pptx') 

            self.menu_clear_layout(self.livres_document_layout)
            self.PPTX_label2 = self.menu_add_label("Vous pouvez modifier le nom du livre",n_cell = 3)
            self.livres_document_layout.addRow(self.PPTX_label2)
            self.OUTFICHIER = self.menu_add_input_text_simple(self.OUTFICHIER)
            self.PPTX_action = self.menu_add_pushbutton("Enregistrer le livre",self.PPTX_save)
            self.livres_document_layout.addRow(self.OUTFICHIER,self.PPTX_action)

            self.livres_album, self.livres_album_layout = self.menu_add_layout("form")
            self.livres_tab.addTab(self.livres_album,"Album")
            self.livres_tab_album()

            self.livres_tab.setCurrentIndex(1)
    #=========================================================================================================
    def PHOTOS_UpdatePhotos_selection(self):
#--------------------------------------------------------------------------------------------------------- 
        format_osxphotos = "yyyy-MM-dd"
        format_edit = "dd-MM-yyyy"

        delta = self.le_écart_photos.currentText()
        if delta == "10 minutes" : d = 600
        elif delta == "30 minutes" : d = 1800


        self.last_date = self.la_date_photos.text()
        str_le_time = self.la_heure_photos.text()
        str_la_date = QDate.toString(QDate.fromString(self.last_date,format_edit),format_osxphotos)
        list_le_date= str_la_date.split("-")

        for idx,photo in enumerate(self.photos_selection):

            if self.le_commentaire_photos.text() != "Sans" and idx == 0 : photo.description = self.le_commentaire_photos.text()

            if self.le_lieu_photos.text() != "Sans": 
                ctx = ssl.create_default_context(cafile=certifi.where())
                geopy.geocoders.options.default_ssl_context = ctx
                geolocator = Nominatim(user_agent="my_geocoder")
                location = geolocator.geocode(self.le_lieu_photos.text())
                if location.latitude and location.longitude : photo.location = (location.latitude,location.longitude)

            str_le_time = QTime.toString(QTime.fromString(str_le_time).addSecs(d))
            list_le_time= str_le_time.split(":")

            photo.date = datetime(int(list_le_date[0]),
                                int(list_le_date[1]),
                                int(list_le_date[2]),
                                int(list_le_time[0]),
                                int(list_le_time[1]),
                                int(list_le_time[2])
                                    )


        self.last_heure = QTime.toString(QTime.fromString(str_le_time).addSecs(60))
        
        self.central_updatePhotos_selection()
    #=========================================================================================================
    def PHOTOS_Import(self):
#--------------------------------------------------------------------------------------------------------- 
        location = None
        if self.le_lieu_photos.text() != "Sans": 
            try: 
                ctx = ssl.create_default_context(cafile=certifi.where())
                geopy.geocoders.options.default_ssl_context = ctx
                geolocator = Nominatim(user_agent="my_geocoder")
                location = geolocator.geocode(self.le_lieu_photos.text())
            except Exception as error: print(error)

        format_osxphotos = "yyyy-MM-dd"
        format_edit = "dd-MM-yyyy"

        delta = self.le_écart_photos.currentText()
        if delta == "10 minutes" : d = 600
        elif delta == "30 minutes" : d = 1800

        self.last_date = self.la_date_photos.text()
        str_le_time = self.la_heure_photos.text()
        str_la_date = QDate.toString(QDate.fromString(self.last_date,format_edit),format_osxphotos)
        list_le_date= str_la_date.split("-")

        import_files = [s for s in os.listdir(import_dir) if os.path.isfile(os.path.join(import_dir, s))]
        import_files.sort(key=lambda s: os.path.getmtime(os.path.join(import_dir, s)))
        import_files_full = [os.path.join(import_dir, s) for s in import_files if s.endswith(".jpg") ]

        photoslib = photoscript.PhotosLibrary()
        import_album = photoslib.album(self.list_selection_albums[-1])
        #imported_photo = photoslib.import_photos([import_files], album=import_album,skip_duplicate_check=True)

        for idx,import_file in enumerate(import_files_full):

            imported_photo_list = photoslib.import_photos([import_file], album=import_album,skip_duplicate_check=True)
            photo = imported_photo_list[0]

        #photoslib.quit()

            if self.le_commentaire_photos.text() != "Sans" and idx == 0 : photo.description = self.le_commentaire_photos.text()
            if location: photo.location = (location.latitude,location.longitude)

            str_le_time = QTime.toString(QTime.fromString(str_le_time).addSecs(d))
            list_le_time= str_le_time.split(":")

            photo.date = datetime(int(list_le_date[0]),
                                int(list_le_date[1]),
                                int(list_le_date[2]),
                                int(list_le_time[0]),
                                int(list_le_time[1]),
                                int(list_le_time[2])
                                    )

        self.last_heure = QTime.toString(QTime.fromString(str_le_time).addSecs(60))

        self.menu_add_infos_row(f'{len(import_files_full)} photo(s) importée(s)')

        for file in os.listdir(import_dir):
            try: os.remove(import_dir+file)
            except Exception as error: print(error)
        
        self.central_import_photos_exec()
    #=========================================================================================================
    def PDF_generate(self):
#--------------------------------------------------------------------------------------------------------- 
        from fpdf import FPDF

        slide_png_dir = self.PDF_leWIP.currentText()
        s_slide_png_dir = slide_png_dir.split("/")
        le_fichier = s_slide_png_dir[-2]
        #path = "/".join(s_slide_png_dir[0:-2])
        output_pdf_path = les_PDF_RV+le_fichier+"_RV.pdf"

        dirFiles = []
        for f in os.listdir(slide_png_dir):
            if ".png" in f :  dirFiles.append(f)

        if dirFiles :
            dirFiles.sort(key=lambda f: int(''.join(filter(str.isdigit,f))))

            fpdf = FPDF(orientation="landscape", format="A4")
            idx_avant_derniere_image = len(dirFiles) - 1
            for idx,f in enumerate(dirFiles):

                if idx == 0: 
                    n_page = PDF_add_page(fpdf,slide_png_dir+f,0)
                    n_page = PDF_add_page(fpdf,"",n_page)

                else :
                    if idx ==idx_avant_derniere_image and n_page%2 == 0 : n_page = PDF_add_page(fpdf,"",n_page)
                    n_page = PDF_add_page(fpdf,slide_png_dir+f,n_page)

            # Save the PDF document
            fpdf.output(output_pdf_path)
            self.infos.append(f'PDF Recto, Verso > {output_pdf_path} généré')
            self.central_accueil()

        else :
            self.menu_add_infos_row("PDF Recto, Verso > Pas de fichier image")

        return
    #=========================================================================================================
    def PPTX_Document_New(self):
#------------------------------------------------------------------------------------------------------------- 
        self.le_document = Presentation(f'{icloud}/MesProgrammes/Mes_Models/{self.model_file.currentText()}.pptx')

        self.PPTX_label.deleteLater()
        self.PPTX_file.deleteLater()
        self.PPTX_action.deleteLater()

        self.PPTX_label2 = self.menu_add_label("Modify PPTX File name",n_cell = 3)
        self.livres_document_layout.addRow(self.PPTX_label2)
        self.OUTFICHIER = self.menu_add_input_text_simple(self.OUTFICHIER)
        self.PPTX_action = self.menu_add_pushbutton("Save PPTX File",self.PPTX_save)
        self.livres_document_layout.addRow(self.OUTFICHIER,self.PPTX_action)

        self.livres_tab.addTab(self.livres_bio, 'Biographies')
        self.livres_tab.addTab(self.livres_album, 'Albums')

        self.livres_tab.setCurrentIndex(1)  
    #=========================================================================================================                       
    def PPTX_publication_album(self):
#-------------------------------------------------------------------------------------------------------------   
#AQS
        le_tri = self.type_tri.currentText()
        isCesureDate = self.CesureDate.isChecked()
        isCesureVille = self.CesureVille.isChecked()
        isCesureComment = self.CesureComment.isChecked()

        avec_personnes = self.avec_personnes.currentText()
        avec_date = self.avec_date.currentText()
        avec_ville = self.avec_ville.currentText()
        avec_fichier = self.avec_fichier.currentText()

        if self.marge_haute.currentText() == "Sans" : 
            y_top = 1
            y_gap = 1
            n_row_max = 2
            x_gap = 1 

        elif self.marge_haute.currentText() == "Petite" : 
            y_top = 20
            y_gap = 2
            n_row_max = 2
            x_gap = 2 
        elif self.marge_haute.currentText() == "Moyenne" : 
            y_top = 40
            y_gap = 1
            n_row_max = 2
            x_gap = 2 
        else:
            y_top = 80
            y_gap = 1
            n_row_max = 2
            x_gap = 2 
#-------------------------------------------------------------------------------------------------------------   
        for idx,le_album_selection in enumerate(self.list_selection_albums):
        #========================

            photos = self.photosdb.photos(albums= [le_album_selection]) 
            register_heif_opener()

            if photos:  
            #=========
                buffer_photos = []
                les_erreurs = []
                for idx,photo in enumerate(photos):
                #==================================
                    if photo.isphoto:  
                    #============                       
                            
                        #"Fichier", "Titre", "Date"
                        if le_tri  == "Fichier":  la_clef = photo.original_filename 

                        elif le_tri  == "Titre": 
                            if photo.title : la_clef = photo.title
                            else: la_clef = photo.date.strftime('%Y%m%d')

                        else : la_clef = photo.date.strftime('%Y%m%d%H%M%S')

                        le_watermark = []
                        # Date
                        le_watermark.append([photo.date.strftime('%d-%m-%Y'),"date"])
                        # Personnes
                        temp_personnes = []
                        for item in photo.persons:
                            if item != "_UNKNOWN_" : temp_personnes.append(item)
                        if temp_personnes : le_watermark.append([temp_personnes,"personne"])
                        #places 
                        if photo.place:
                            le_texte = ""  
                            if photo.place.names.area_of_interest : 
                                le_texte += f"{photo.place.names.area_of_interest[0]}"
                                le_watermark.append([photo.place.names.area_of_interest[0],"interest"])
                            if photo.place.names.street_address : 
                                le_texte += f" {photo.place.names.street_address[0]}"
                                le_watermark.append([photo.place.names.street_address[0],"street"])
                            if photo.place.names.city : 
                                le_texte += f" {photo.place.names.city[0]}"
                                le_watermark.append([photo.place.names.city[0],"city"])
                            if photo.place.names.country : 
                                le_texte += f" {photo.place.names.country[0]}"
                                le_watermark.append([photo.place.names.country[0],"country"])
                            if le_texte : le_watermark.append([le_texte,"place"])  
                        else:
                            le_watermark.append(["non défini","city"])
                        # Commentaire
                        if photo.description:  le_watermark.append([photo.description,"description"])

                        le_watermark.append([photo.original_filename ,"fichier"])

                        if photo.title : le_watermark.append([photo.title,"titre"])

                        hasRgb_im = False
                        for filephoto in photo.path_derivatives:
                            rgb_im = cv2.imread(filephoto)
                            try:  
                                len(rgb_im)
                                hasRgb_im = True
                                break
                            except : next

                        if hasRgb_im : 
                            image_height, image_width = rgb_im.shape[:2]
                            photo_filename = f'{export_path}/{idx}.jpg'
                            cv2.imwrite(photo_filename, rgb_im)
                            buffer_photos.append([la_clef,photo_filename,photo.width,photo.height,le_watermark,photo.favorite,image_width,image_height])
                        else:
                            print(f"{photo.original_filename} problème path_derivatives")
                    else: les_erreurs.append(f"{photo.original_filename} n'est pas une photo")
                        
                if les_erreurs:
                    for item in les_erreurs:
                        print(item)
# Add photos to document
#-----------------------------------------------------------------------------------------------------------------------------------------
                if buffer_photos:
                    buffer_photos = sorted(buffer_photos, key=lambda col: (col[0]) )
#page de garde 
                    la_page_garde,boxes = PPTX_add_page(self.le_document,slide_layout_garde)

                    temp_texte = le_album_selection.replace("LIVRE_","")
                    t_zorg = temp_texte.split("@")
                    if len(t_zorg) > 1: temp_texte_place= t_zorg[-1]
                    else: temp_texte_place = ""
                    temp_texte_date = ""
                    temp_texte_date = t_zorg[0].split(" ")[0]
                    temp_texte = " ".join(t_zorg[0].split(" ")[1:])

                    p = PPTX_add_paragraph(boxes[0])
                    if temp_texte : PPTX_add_run(p,temp_texte)
                    if temp_texte_date : 
                        p = PPTX_add_paragraph(boxes[0])
                        PPTX_add_run(p,temp_texte_date)
                    if temp_texte_place:
                        p = PPTX_add_paragraph(boxes[0])
                        PPTX_add_run(p,temp_texte_place)

# initialisation des dimensions

                    n_row = 1
                    x_bandeau = x_gap
                        
                    x = slide_width - x_gap
                    y = y_top

                    la_page = False
                    isPreviousFavorite = False
                    isNext_row = False
                
                    la_date_previous = ""
                    la_ville_previous = ""
                    le_commentaire_previous = ""
                    le_commentaire_cur = ""
                    pic_previous = None

                    la_ville_cur = ""
                    #photo_height = min ((slide_height - y_top)/n_row_max - y_gap, 194/n_row_max)
                    photo_height = (slide_height - y_top)/n_row_max - y_gap
                    #print(photo_height,(slide_height - y_top)/n_row_max - y_gap)

        # boucle sur les phto du buffer photos
                    isNewPage = False
                    le_texte_p0 =""
                    for idx,buffer_photo in enumerate(buffer_photos):

                        y_diff = 0

# initiage image from local file
                        le_fichier_image = buffer_photo[1]
# les watermarks
                        for item in buffer_photo[4]:
                            if item[1] == "date" : la_date_cur = item[0]
                            if item[1] == "city" : la_ville_cur = item[0]
                            if item[1] == "description" : le_commentaire_cur = item[0]
# resize image
                        percent_crop_horizontal = 0.0
                        percent_crop_vertical = 0.0

                        r = buffer_photo[6] / buffer_photo[7]
                        img_height = photo_height
                        img_width = photo_height * r 

                        if isCesureDate  :
                            if la_date_cur != la_date_previous:
                                la_date_previous = la_date_cur
                                isNewPage = True

                        if isCesureVille:
                            if la_ville_cur != la_ville_previous:
                                la_ville_previous = la_ville_cur
                                isNewPage = True

                        if isCesureComment:
                            if le_commentaire_cur != le_commentaire_previous:
                                le_commentaire_previous = le_commentaire_cur
                                le_texte_p0 = le_texte_p0 + le_commentaire_cur
                                isNewPage = True                           
                    
                        elif not la_page:
                            isNewPage = True
# photo favorite IICI
                        if buffer_photo[5]: 
                            iw = buffer_photo[6]
                            ih = buffer_photo[7]
                            if iw > ih : 
                                img_height = slide_height - y_top -y_gap
                                img_width = img_height * r
                                x_diff =  img_width - slide_width + 2 * x_gap
                                if x_diff > 0: percent_crop_horizontal = x_diff/img_width                                 
                            else:
                                img_width = slide_width  - 2 * x_gap
                                img_height = img_width / r 
                                y_diff =  img_height - slide_height + y_top + y_gap
                                if y_diff > 0: percent_crop_vertical = -y_diff/img_height   

                            isNewPage = True  
                            isPreviousFavorite = True 
                        elif  isPreviousFavorite:
                            isPreviousFavorite = False
                            isNewPage = True
# Nouvelle page
                        if isNewPage:

                            x = x_gap
                            y = y_top
                            
                            la_page, boxes = PPTX_add_page(self.le_document,slide_layout_album)
                            p0 = PPTX_add_paragraph(boxes[0])
                            if isCesureDate or isCesureVille : 
                                PPTX_add_run(p0,f'{la_date_cur}, {la_ville_cur}',"bold")
                                p0 = PPTX_add_paragraph(boxes[0])

                            isNewPage = False
                            isNext_row = False
                            n_row = 1
                            
                        else:

                            x_diff = int( slide_width - x  - img_width - x_bandeau)

                            if x_diff  < 0 : 

                                if slide_width - x - x_gap - x_bandeau > 30 :
                                    percent_crop_horizontal = -1 * (x_diff / img_width)
                                    isNext_row = True
                                else:
                                    isNext_row = False
                                    PPTX_expand_picture(pic_previous,x_bandeau)
                                    n_row = n_row + 1
                                    
                                    if n_row > 2 :
# New page
                                        x = x_gap
                                        y = y_top
                                        
                                        la_page, boxes = PPTX_add_page(self.le_document,slide_layout_album)
                                        p0 = PPTX_add_paragraph(boxes[0])
                                        if isCesureDate or isCesureVille : 
                                            PPTX_add_run(p0,f'{la_date_cur}, {la_ville_cur}',"bold")
                                            p0 = PPTX_add_paragraph(boxes[0])

                                        isNewPage = False
                                        isNext_row = False
                                        n_row = 1

                                    else :
                                        x = x_gap
                                        y = y + img_height + y_gap                         
# Ajout de l'image   
                        xp = x

                        if img_width > 0 and img_height > 0 :
                            pic = la_page.shapes.add_picture(le_fichier_image,Mm( xp ),Mm(y - y_diff),Mm(img_width),Mm(img_height))
                            pic.crop_left = percent_crop_horizontal
                            pic.crop_right = -1 * percent_crop_horizontal
                            pic.crop_top = percent_crop_vertical
                            pic.crop_bottom = -1 * percent_crop_vertical
                            pic_previous = pic

                        else:
                            print("img_width=",img_width, "img_height=",img_height)
                            exit()
# Ajout du texte
                        if img_width > 4:
# Ajout du textex
                            les_w_textes = []
                            les_personnes = []
                            les_comments = []

                            for item in buffer_photo[4]:
                                
                                if item[1] == "description":
                                    if item[0] not in les_comments:
                                        PPTX_add_run(p0,f'{item[0]}. ')
                                        les_comments.append(item[0])

                                if ( 
                                    (avec_ville == "Avec" and item[1] == "city" and item[0] != "non défini")
                                or (avec_date == "Avec" and item[1] == "date")
                                or (avec_fichier == "Avec" and item[1] == "fichier")
                                ):
                                    les_w_textes.append(item[0])

                                if avec_personnes == "Avec" and item[1] == "personne": les_personnes = item[0]

                            les_w_textes = les_w_textes + les_personnes
                            if les_w_textes:
                                #object_textbox = PPTX_add_box(la_page,xp+2,y+2,img_width - 4 , 6.65 * len(les_w_textes),
                                object_textbox = PPTX_add_box(la_page,xp+2,y+2,img_width/2 , 6.65 * len(les_w_textes),
                                                                "transparent",bcolor=Gray1)

                                for item in les_w_textes:
                                    p = PPTX_add_paragraph(object_textbox)
                                    PPTX_add_run(p,item,size = 14,color = White) 

# page Recto
                        if isNext_row :
                                n_row = n_row + 1
                                
                                if n_row > 2 :
                                    isNewPage = True
                                else :
                                    x = x_gap
                                    y = y + photo_height + y_gap  
                                    isNext_row = False
                        else:

                            x = x + img_width + x_gap 
#-------------------------------------------------------------------------------------------------------------
                    try:
                        for filename in glob.iglob(f"{export_path}/*.*", recursive=False):
                            os.remove(filename)
                    except:
                        print("Erreur vidage de répertoire")
#-------------------------------------------------------------------------------------------------------------
                    self.livres_tab.setCurrentIndex(0)
                    self.les_infos = f"+ Album {le_album_selection} {le_tri}"
                    self.menu_add_infos_row("Albums > "+self.les_infos)

                    self.livres_tab_album()

#==============================================================================================================   
#==============================================================================================================                           
#==============================================================================================================   
if __name__ == "__main__":
    import sys
    app = QApplication(sys.argv)
    MainWindow = LA_WINDOW()
    MainWindow.show()
    sys.exit(app.exec())
#==============================================================================================================   
#==============================================================================================================                           
#==============================================================================================================   
