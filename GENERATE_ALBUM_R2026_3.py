import os
import psutil
# file management
import glob
import os.path
from datetime import datetime
import subprocess
import photoscript
import geopy,certifi,ssl
from geopy.geocoders import Nominatim

#= Image Management ============================================
from PIL import Image, ImageDraw, ImageFont
from pillow_heif import register_heif_opener
import cv2
#===============================================================
#-- _PPTX_
from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.enum.shapes import *
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import *
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement

#============================================================================================================
class my_photo(object):
    def __init__(self):
        p_def = {
        "id" : None,
        "file_original" : None,
        "filename" : None, 
        "date" : None,
        "datelong" : None,
        "description" : None,
        "title" : "",
        "width" : 0,
        "height" : 0,
        "favorite" : False,
        "person" : None,
        "poi" : None,
        "street" : None,
        "city" : None,
        "country" : None,
        "place" : None}
        for key, value in p_def.items():
            setattr(self, key, value)

    def __eq__(self, other):
        if list(self.__dict__.keys()) == list(other.__dict__.keys()): #print("same dict")
            if (list(self.__dict__.values())) == (list(other.__dict__.values())): return True #print("same value")
            else : return False #print("different values")
        else: return False #print("different dict")

#============================================================================================================
#
page_width = 190
pt_mm = 0.68
document_font = "Aptos (Corps)"
font_file = "/Users/bernardconti/Library/Fonts/Aptos.ttf"
font_file_bold = "/Users/bernardconti/Library/Fonts/Aptos-Bold.ttf"
# --------------------------------------------------------------
icloud = "/Users/bernardconti/Library/Mobile Documents/com~apple~CloudDocs"
les_WIP = "/Users/bernardconti/LOCAL_TEMP/WIP/"
les_PDF_RV = "/Users/bernardconti/Library/Mobile Documents/com~apple~CloudDocs/Les éditions du 57/PDF_RectoVerso/"
icon_dir = icloud+'/MesProgrammes/python_global/icons'
export_path =  "/Users/bernardconti/LOCAL_TEMP/Photos"
save_path = "/Users/bernardconti/LOCAL_TEMP/Documents"
import_dir = "/Users/bernardconti/LOCAL_TEMP/ePhoto_Doc/Import/"
watermarkdir = "/Users/bernardconti/LOCAL_TEMP/watermark"

#--------------------------------------------------------------
retour_ligne = '\n'
Gray= "#EFEEEF"
GrayRow= "#EFEEEF"
Black ="#000000"
White = "#FFFFFF"
couleur_homme ="#ffcc99"
couleur_femme="#ccccff"
font_texte ="Aptos (Corps)"
couleur_cible = "#30DD30"
Gray1= "#b7b1b1"#F2F2F2
GraySide = "#c5c5c5"
Gray3 ="#D9D6D6"
Gray4 = "#E1DEDEEF"
couleur_chemin = "#EC5800"
couleur_Title ="#EB742F"
Green = "#B4E1C0"
Gray2= "#fde9d9"
Silver = "#A9A6A6" 

#Bleu = "#375e94"
Bleu = "#272170"
Darkblue = "#102A37"
Gris_clair = "#EFEFEF"
ligne = "#375e94"
descendant = "#000000"
red = "#FF0000"
page_width = 190
pt_mm = 0.68
document_font = "Aptos (Corps)"
font_file = "/Users/bernardconti/Library/Fonts/Aptos.ttf"
font_file_bold = "/Users/bernardconti/Library/Fonts/Aptos-Bold.ttf"

EMU = 36000

#PPT
slide_width = 277
slide_height = 210
slide_margin_left = 5
slide_margin_right = 5
slide_margin_top = 10
slide_margin_bottom = 5

le_model = icloud+"/MesProgrammes/python_Albums/MODEL_ALBUM.pptx"
slide_layout_album_break = 0
slide_layout_album_breakfree = 1
slide_layout_garde = 2
slide_layout_57 = 3

le_main_folder = "LES LIVRES_CONTI"
#=============================================================================================================
# BOXES
#=============================================================================================================
def PPTX_add_box(la_page,x,y,w,h,*args,**kwargs):
#------------------------------------------------------------------------------------------------------------- 
    if w < 0 or h < 0 or not la_page :  
        print("PPTX_add_box",w,h)
        return False
#
# lecture des parametres
#
    la_couleur = False
    le_alignement = MSO_ANCHOR.MIDDLE
    la_margin_bottom = 0.5
    la_margin_top = 0.5
    la_margin_left = 0.5
    la_margin_right = 0.5
        
    for clef, valeur in kwargs.items():  
        if valeur: 
            if isinstance(clef,str): clef =clef.lower()
            if clef == "bcolor" : la_couleur = valeur
            if clef == "margin_bottom" : la_margin_bottom = valeur
            if clef == "margin_top" : la_margin_top = valeur
            if clef == "margin_left" : la_margin_left = valeur
            if clef == "margin_right" : la_margin_right = valeur

    isTransparent = False
    isWrap = True
    isAutosize = False
    isCadre = False

    for valeur in args:
        if valeur:
            if isinstance(valeur, str): valeur =valeur.lower()
            if valeur == "nowrap" : isWrap =  False    
            if valeur == "transparent" : isTransparent =  True  
            if valeur == "autosize" : isAutosize =  True    
            if valeur == "cadre" : isCadre =  True 

    # ajout box
    txBox = la_page.shapes.add_textbox(Mm(x),Mm(y),Mm(w),Mm(h))

    if isCadre : 
        line = txBox.line
        line.color.rgb = RGBColor(200,200,200)
    elif la_couleur:
            txBox.fill.solid()
            txBox.fill.fore_color.rgb = hex_to_rgb(la_couleur)

    if isTransparent:
                            
        def SubElement(parent, tagname, **kwargs):
            element = OxmlElement(tagname)
            element.attrib.update(kwargs)
            parent.append(element)
            return element

        def _set_shape_transparency(shape, alpha):
        # Set the transparency (alpha) of a shape
            ts = shape.fill._xPr.solidFill
            sF = ts.get_or_change_to_srgbClr()
            sE = SubElement(sF, 'a:alpha', val=str(alpha))
#ZOB
        _set_shape_transparency(txBox,66000)

    tf = txBox.text_frame
    tf.word_wrap = isWrap
    tf.auto_size = isAutosize 
    if le_alignement : tf.vertical_anchor = le_alignement
    tf.margin_bottom = Mm(la_margin_bottom)
    tf.margin_top = Mm(la_margin_top)
    tf.margin_left = Mm(la_margin_left)
    tf.margin_right = Mm(la_margin_right)
#-------------------------------------------------------------------------------------------------------------
    return txBox
#=============================================================================================================
def PPTX_add_paragraph(text_Box,*args,**kwargs):
#------------------------------------------------------------------------------------------------------------- 
    align = PP_ALIGN.LEFT
    level = 0
    before_space = 0

    for valeur in args:
        if valeur:
            if isinstance(valeur, str): valeur =valeur.lower()
            if valeur == "center" : align =  PP_ALIGN.CENTER    
            if valeur == "right" : align =  PP_ALIGN.RIGHT           

    for clef, valeur in kwargs.items():  
        if valeur: 
            if isinstance(valeur,str): clef = clef.lower()
            if clef == "level" : level = int(valeur)
            if clef == "before_space" : before_space = int(valeur)
#-------------------------------------------------------------------------------------------------------------  
    tf = text_Box.text_frame
    if len(tf.paragraphs[0].text) == 0 : le_paragraph = tf.paragraphs[0] 
    else: le_paragraph = tf.add_paragraph()

    le_paragraph.level = level 
    le_paragraph.alignment = align
    le_paragraph.space_before = Pt(before_space)

    return le_paragraph
#=============================================================================================================
def PPTX_add_run(le_paragraph,le_texte,*args,**kwargs):
#------------------------------------------------------------------------------------------------------------- 
    run = None
    if le_texte:
#------------------------------------------------------------------------------------------------------------- 
        box_width = False
        isItalic =False
        isBold = False
        isUnderline = False
        font = False
        size = False
        color = False
        
        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "box_width" : box_width = valeur
                if clef == "size" : size = int(valeur)
                if clef == "font" : font = valeur
                if clef == "color" : color = valeur

        for valeur in args:
            if valeur:
                if isinstance(valeur, str): valeur =valeur.lower()
                if valeur == "bold" : isBold =  True    
                if valeur == "italic" : isItalic =  True  
                if valeur == "underline" : isUnderline =  True    
#-------------------------------------------------------------------------------------------------------------   
        if box_width:
            w_chr = (box_width-1)/len(le_texte)
            size = w_chr * 5
            if size > 12.5:size = 12.5
            if size < 8: size = 8
#------------------------------------------------------------------------------------------------------------- 
        run = le_paragraph.add_run()
        run.text =le_texte
        if font : run.font.name = font
        if size : run.font.size = Pt(size)
        if color : run.font.color.rgb = hex_to_rgb(color)        
        if isBold : run.font.bold = isBold
        if isItalic : run.font.italic = isItalic
        if isUnderline : run.font.underline = isUnderline

    return run
#=============================================================================================================
def PPTX_add_page(le_document,le_layout):
    boxes = []
    la_page = le_document.slides.add_slide(le_document.slide_layouts[le_layout])
    for shape in la_page.placeholders:
        #print('%d %s' % (shape.placeholder_format.idx, shape.name))
        boxes.append(shape)
    b = PPTX_add_box(la_page,slide_width-10,slide_height-5,10,5)
    n_slide = le_document.slides.index(la_page)
    if n_slide > 1 : PPTX_add_run(PPTX_add_paragraph(b,"right"),str(n_slide),"italic","bold",size=10)
    return la_page,boxes
#=============================================================================================================
def hex_to_rgb(value):
#-------------------------------------------------------------------------------------------------------------
    RGB = RGBColor(0, 0, 0)
    try:
        value = value.lstrip('#')
        lv = len(value)
        (rouge, vert, bleu) = tuple(int(value[i:i + lv // 3], 16) for i in range(0, lv, lv // 3))
        RGB = RGBColor(rouge, vert, bleu)
    except Exception as error:
        print (value,tuple(int(value[i:i + lv // 3], 16) for i in range(0, lv, lv // 3)))
#-------------------------------------------------------------------------------------------------------------
    return RGB
#=============================================================================================================
def box_height(le_content,la_width,le_style,la_fontsize,top_margin,bottom_margin,left_margin,right_margin):

    h_row = 0
    # gestion des tables
    la_width = la_width - left_margin - right_margin

    if le_style == "bold" : la_font = ImageFont.truetype(font_file_bold,la_fontsize)
    else:la_font = ImageFont.truetype(font_file,la_fontsize)

    ascent, descent = la_font.getmetrics()

    le_h = 0
    if le_content:

        les_textes = []
        if isinstance(le_content,str):      les_textes = [le_content]
        elif isinstance(le_content,list):   les_textes = le_content
        elif isinstance(le_content,int):    les_textes = [str(le_content)]

        if les_textes : 
            for le_texte in les_textes:
                les_textes_sp= le_texte.splitlines()
                for item in les_textes_sp:
                    w = la_font.getlength(item) * 0.352778
                    n = int(w/(la_width))+1
                    le_h = le_h + (ascent + descent)*0.352778*1*n

            le_h = le_h + top_margin + bottom_margin
            

    return le_h
#=============================================================================================================
def traduction_date(texte):
    texte = texte.lower()
    texte = texte.replace("january","janvier")
    texte = texte.replace("february","février")
    texte = texte.replace("april","avril")
    texte = texte.replace("may","mai")
    texte = texte.replace("june","juin")
    texte = texte.replace("july","juillet")
    texte = texte.replace("august","août")
    texte = texte.replace("september","septembre")
    texte = texte.replace("october","octobre")
    texte = texte.replace("novembre","novembre")
    texte = texte.replace("december","décembre")
    texte = texte.replace("monday","Lundi")
    texte = texte.replace("tuesday","Mardi")
    texte = texte.replace("wednesday","Mercredi")
    texte = texte.replace("thursday","Jeudi")
    texte = texte.replace("friday","Vendredi")
    texte = texte.replace("saturday","Samedi")
    texte = texte.replace("sunday","Dimanche")
    return texte


#====PDF =====================================================================================================
def PDF_add_page(le_pdf,image_filename,n_page):
#=============================================================================================================
    n_page = n_page + 1
    le_pdf.add_page()
    
    if n_page % 2 == 0: 
        xi = 0
        xt = 282
        xc = 287
    else: 
        xi = 20
        xt = 5
        xc = 10
    
    if n_page > 1:
        le_pdf.set_font('helvetica', size=12)
        le_pdf.text(x= xt ,y=205,text=f'# {n_page}')
    else:        
        le_pdf.set_font('helvetica', size=10)
        le_pdf.text(x= xt ,y=200,text=f'Édition')
        le_pdf.text(x= xt ,y=205,text=f'du 57')

    le_pdf.circle(xc,55,2.5,"F")
    le_pdf.circle(xc,155,2.5,"F")

    if image_filename :  le_pdf.image(image_filename,  x = xi , y = 0 ,w= 277,h=210)

    return n_page
#=============================================================================================================
def excute_cmd(command): 
# --------------------------------------------------------------   
    command_list = command.split(" ")
    #print(command_list)
    result = subprocess.run(command_list, capture_output=True, text=True)
    if result.stdout : print ("resultat = " + result.stdout)
    if result.stderr : print("erreur = " + result.stderr)
# --------------------------------------------------------------
    return(result.stdout,result.stderr)
#=============================================================================================================

#=============================================================================================================
def PPTX_expand_picture(pic,x_bandeau) :
        if pic:
            pvh = pic.height
            r = pic.width / pvh

            pic.width = (slide_width - x_bandeau ) * EMU - pic.left 
            pic.height = int(pic.width / r)

            pic.crop_bottom = -(1-pvh/pic.height)
            pic.crop_top = (1 - pvh/pic.height)
        return
#=============================================================================================================

#===========================================================================================
# DEBUT
#===========================================================================================
print("Execution de ",__file__)
dir  = "/Users/bernardconti/Downloads/"

#-- MENU ---------------------------------------------------------------------------------------
from PyQt6.QtWidgets import *
from PyQt6.QtGui import *
from PyQt6.QtCore import *

custom_font = QFont("Arial",14)
line_height = 30
line_spacing = 10
cell_spacing = 5
n_line = 7
bio_width = [400,120,120,120,120,120]
bio_layout_width = 0
for number in bio_width:
    bio_layout_width += number
bio_layout_width += (len(bio_width)- 1)*cell_spacing
bio_layout_height = n_line * line_height + (n_line-1)*line_spacing

cell_width = 250

#========================================================================================    
for proc in psutil.process_iter():
    if proc.name() == "Microsoft PowerPoint": proc.kill()
#========================================================================================
# recupere tous les noms
#========================================================================================

# initiate list des directore pour impression PDF R&V
list_dirs = glob.glob(les_WIP+"*/", recursive = False)
if not list_dirs : exit()
#========================================================================================   
#  MAIN WINDOW   
#========================================================================================
class LA_WINDOW(QMainWindow):
#========================================================================================
    def __init__(self):
        #super().__init__(*args, **kwargs)
        super().__init__()

        import osxphotos 
        try:
            self.photosdb = osxphotos.PhotosDB()
            #self.photosdb = osxphotos.PhotosDB(dbfile = "/Users/bernardconti/Pictures/PhotosLibrary_BC.photoslibrary")
        except Exception as error:
            print("Photos DB",error)
            exit()

        self.setupUi()
#----------------------------------------------------------------------------------------
    def central_accueil(self):
        self.menu_clear_layout(self.central_layout)
        for item in self.infos:
            info = QLabel(item)
            info.setFixedSize(bio_layout_width, line_height)
            info.setStyleSheet("background-color: lightgray;color: black;"
                                    "qproperty-alignment: AlignLeft;"
                                    "border-radius: 10px;"
                                    "qproperty-wordWrap: true;"
                                    "padding: 6px;"
                                    )
            info.setFont(custom_font)
            self.central_layout.addRow(info) 
        self.menu_add_infos_row("2026, by Bernard CONTI, Dourdan, France")
#===========================================================================================================
    def setupUi(self):

        self.setObjectName("Main Menu")
        self.resize(400, 400)

# Create menu bar, toolbar and statusbar objects
        menubar = self.menuBar()
        menubar.setFont(custom_font)
        
        self.centralwidget = QWidget(self)
        self.setCentralWidget(self.centralwidget)
        self.statusbar = QStatusBar(self)

# Create actions
        finAction = self.menu_add_action(icon = "exit.png",text = "Good Bye !!!", cmd = self.PPTX_fin)
        #updatePhotos_album_Action = self.menu_add_action(text = "Update Photos Album", cmd = self.central_updatePhotos_album)
        importPhotos_Action = self.menu_add_action(text = "Import Photos", cmd = self.central_import_photos)
        updatePhotos_selection_Action = self.menu_add_action(text = "Update Photos Selection", cmd = self.central_updatePhotos_selection)
        negatifsAction = self.menu_add_action(text = "Négatifs", cmd = self.central_negatifs)

        pdfAction = self.menu_add_action(text = "PDF Recto Verso", cmd = self.central_pdf)
        livresAction = self.menu_add_action(text = "Album", cmd = self.central_livres)
        BienvenueAction = self.menu_add_action(icon = "home.png", text = "Home Sweet Home",cmd = self.central_accueil)
# Creation des mennus     
        helpMenu = menubar.addMenu(QIcon(icon_dir+"/tools.png"),"&Tool Box")
        helpMenu.setStyleSheet("background-color: "+ couleur_homme +";color: black;"
        "font-name=Arial;"
        "font-size: 16px;"
        )
# Add actions to menus
        #helpMenu.addAction(updatePhotos_album_Action)
        self.last_date = None
        self.last_heure = None
        helpMenu.addAction(importPhotos_Action)
        helpMenu.addAction(updatePhotos_selection_Action)
        helpMenu.addAction(negatifsAction)
        helpMenu.addAction(pdfAction)

# add menus to menubar
        menubar.addMenu(helpMenu)

# add  toolbar
        toolbar = QToolBar('Main ToolBar', self)
        toolbar.setIconSize(QSize(24, 24))
        toolbar.setStyleSheet("background-color: "+ couleur_femme +";color: black;"
                    "font-name=Arial;"
                    "font-size: 16px;"
                    )

        toolbar.addAction(BienvenueAction)
        toolbar.addAction(livresAction)
        toolbar.addActions(menubar.actions())
        toolbar.addAction(finAction)
        
# add toolbar and statusbar to main window
        self.addToolBar(toolbar)
        self.setStatusBar(self.statusbar)

# central_layout, self
        self.central_layout = QFormLayout(self.centralwidget) 
        self.central_layout.setVerticalSpacing(line_spacing)
        self.central_layout.setHorizontalSpacing(cell_spacing) 
#init central on welcome
        
        self.infos = ["Welcome to the ALBUM generator"]
        self.central_accueil()

#======================================================================================================================          
    def PPTX_save(self):
        la_page = self.le_document.slides.add_slide(self.le_document.slide_layouts[slide_layout_57])
        OUT_FileName  = f'{save_path}/{self.OUTFileName}'
        self.le_document.save(OUT_FileName)

        for file in os.listdir(watermarkdir):
            try: os.remove(watermarkdir+"/"+file)
            except Exception as error: print(error)

        self.infos.append("Album > "+ OUT_FileName +  " saved")
        self.central_accueil()
#==============================================================================================================            
    def PPTX_fin(self):
        print("Good Bye !!!")
        self.force_close = True
        self.close()
#==============================================================================================================

#==============================================================================================================
# DEF_menu_add_...
#============================================================================================================== 
    def list_albums(self,le_layout,cmd):
#==============================================================================================================
        self.menu_clear_layout(le_layout)

        self.next = self.menu_add_pushbutton('Suite',cmd,n_cell = 2)

        self.tree = QTreeWidget()
        self.tree.setHeaderHidden(True)
        self.tree.setFixedSize(2*cell_width, 10* line_height)

        selmodel = self.tree.selectionModel()
        selmodel.selectionChanged.connect(self.click_action)

        le_layout.addRow(self.tree)
        le_layout.addRow(self.next)

        self.list_selection_albums = []
        self.les_folders = []
        for f in self.photosdb.folder_info:
            if f.title == le_main_folder: 
                for fs in f.subfolders:

                    isP = True
                    for le_album in fs.album_info:
                        if le_album.title[0] != "P":
                            isP = False
                            break

                    if not isP:
                        self.les_folders.append([fs.title,fs])
                        level0 = QTreeWidgetItem(self.tree, [fs.title])
                        for le_album in fs.album_info:
                            if le_album.title[0] != "P": level1 = QTreeWidgetItem(level0, [le_album.title])                 
#===========================================================================================
    def click_action(self, selected, deselected):

        for index in selected.indexes():
            item = self.tree.itemFromIndex(index)
            if self.tree.itemFromIndex(index.parent()) : 
                if item.text(0) not in self.list_selection_albums : 
                    self.list_selection_albums.append(item.text(0))
                    self.menu_add_infos_row(f'{self.tree.itemFromIndex(index.parent()).text(0)} > {item.text(0)} ')
            else:
                for f in self.les_folders:
                    if item.text(0) == f[0]:
                        for le_album in f[1].album_info:
                            if le_album.title not in self.list_selection_albums : 
                                if le_album.title[0] != "P": self.list_selection_albums.append(le_album.title)
                        self.menu_add_infos_row(f'Dossier : {item.text(0)} ')
#================================================================================================
    def livres_tab_album_exec(self):

        if not self.list_selection_albums : 
            print("No Folder or Album selected")
            return

        self.menu_clear_layout(self.central_layout)

        self.central_layout.addRow(self.menu_add_label("Albums",n_cell = 2,bcolor = "green",color = "white"))
        #for a in self.list_selection_albums:
            #self.central_layout.addRow(self.menu_add_label(a,n_cell = 2))
        le_FileName = f'ALBUM_{"|".join(self.list_selection_albums).replace("LIVRE_","")}.pptx'
        le_FileName = le_FileName.replace(" ","")
        self.OUTFileName = le_FileName
        self.central_layout.addRow(self.menu_add_label(self.OUTFileName,n_cell = 2))
# section regroupement      
        self.central_layout.addRow(self.menu_add_label("Sort type :",n_cell = 2,bcolor = "green",color = "white"))
# type de tri
        self.sorttype  = self.menu_add_input_text_combobox(["Date","Title","FileName"],n_cell = 1)
        self.central_layout.addRow(self.menu_add_label("Photo Sorted by"),self.sorttype)
# Changement de slide
        self.is_breakdate = self.menu_add_checkbox("New date > New page",n_cell = 1)
        self.central_layout.addRow(self.menu_add_label(" ",bcolor=White),self.is_breakdate)
# Changement de slide
        self.is_breakcity = self.menu_add_checkbox("New city > New page",n_cell = 1)
        self.central_layout.addRow(self.menu_add_label(" ",bcolor=White),self.is_breakcity)
# Changement de slide
        self.is_breakdescription = self.menu_add_checkbox("New comment > New Page",n_cell = 1)
        self.central_layout.addRow(self.menu_add_label(" ",bcolor=White),self.is_breakdescription)
# section Marge Haute      
        self.central_layout.addRow(self.menu_add_label("Photos disposal :",n_cell = 2,bcolor = "green",color = "white"))
# Marge haute
        self.marge_haute = self.menu_add_input_text_combobox(["Compact","Dense","Classic"],n_cell = 1)
        self.central_layout.addRow(self.menu_add_label("Marge Haute"),self.marge_haute)

        # section info_box      
        self.central_layout.addRow(self.menu_add_label("Info Box Configuration :",n_cell = 2,bcolor = "green",color = "white"))
        # Info Persons    
        self.is_person = self.menu_add_checkbox("Person",n_cell = 1)
        self.central_layout.addRow(self.menu_add_label(" ",bcolor=White),self.is_person)
        # Info Dates     
        self.is_date = self.menu_add_checkbox("Date",n_cell = 1)
        self.central_layout.addRow(self.menu_add_label(" ",bcolor=White),self.is_date)
        # Info City  
        self.is_city = self.menu_add_checkbox("City",n_cell = 1)
        self.central_layout.addRow(self.menu_add_label(" ",bcolor=White),self.is_city)
        # Info description 
        self.is_description = self.menu_add_checkbox("Description",n_cell = 1)
        self.central_layout.addRow(self.menu_add_label(" ",bcolor=White),self.is_description)
        # Info Title 
        self.is_title = self.menu_add_checkbox("Title",n_cell = 1)
        self.central_layout.addRow(self.menu_add_label(" ",bcolor=White),self.is_title)
        # Info FileName  
        self.is_filename = self.menu_add_checkbox("FileName",n_cell = 1)
        self.central_layout.addRow(self.menu_add_label(" ",bcolor=White),self.is_filename)

        #----------------------------------------------------------------------------------------------------            
        # Action de Publication
        self.central_layout.addRow(self.menu_add_pushbutton("Publish Albums",self.PPTX_publication_album))      
#==============================================================================================================
    def central_livres(self):
        self.menu_clear_layout(self.central_layout)
        self.list_albums(self.central_layout,self.livres_tab_album_exec)
#===========================================================================================================
    def central_pdf(self):
        self.menu_clear_layout(self.central_layout)
        self.PDF_leWIP = self.menu_add_input_text_combobox(list_dirs,n_cell = 2)
        self.pdf_action = self.menu_add_pushbutton("Générer le PDF Recto Verso",self.PDF_generate)
        self.central_layout.addRow(self.PDF_leWIP,self.pdf_action)
#===========================================================================================================           
    def central_negatifs(self):
        self.menu_clear_layout(self.central_layout)       
        
        save_path = "/Users/bernardconti/LOCAL_TEMP/Négatifs/"
        target_album="NÉGATIFS_CONVERTIS"

        for idx,photo in enumerate(self.photosdb.photos(albums=["NÉGATIFS_TEMP"])):

            image = cv2.imread(photo.path_derivatives[0])
            gray_image = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
            negative_image = 255 - gray_image 
            ct=datetime.now()
            savefilename = f'{save_path}NEGATIF_{ct.year}-{ct.month}-{ct.day}_{ct.hour}:{ct.minute}:{ct.second}_{idx}.jpg'       
            self.central_layout.addRow(self.menu_add_label(f'NEGATIF_{ct.year}-{ct.month}-{ct.day}_{ct.hour}:{ct.minute}:{ct.second}_{idx}.jpg',n_cell = 2 ))
            cv2.imwrite(savefilename, negative_image) 
            command_list = ["osxphotos","import",savefilename,"--album",target_album,"--skip-dups"]
            result = subprocess.run(command_list, capture_output=False, text=True)
            os.remove(savefilename)
#===========================================================================================================
    def central_updatePhotos_selection(self):
        self.menu_clear_layout(self.central_layout) 
        self.central_layout.addRow(self.menu_add_pushbutton("Sélection des photos",self.central_updatePhotos_selection_exec, n_cell = 2, bcolor = "green", color = "white"))
#===========================================================================================================
    def central_updatePhotos_selection_exec(self):

        photoslib = photoscript.PhotosLibrary()
        if photoslib.selection:
            self.photos_selection = photoslib.selection
            for la_photo in self.photos_selection:
                self.central_layout.addRow(self.menu_add_label(f'{la_photo.filename}',n_cell = 2))

        self.la_date_photos = self.menu_add_date(default = self.last_date,ncol = 2)
        self.la_heure_photos = self.menu_add_heure(default = self.last_heure,ncol = 2)
        self.le_écart_photos =self.menu_add_input_text_combobox(["10 minutes", "30 minutes", "1 heure"],n_cell = 1)
        self.le_commentaire_photos =self.menu_add_input_text_simple("No",n_cell = 1)
        self.le_lieu_photos = self.menu_add_input_text_simple("No",n_cell = 1)

        self.central_layout.addRow(self.menu_add_label("Jour"),self.la_date_photos)
        self.central_layout.addRow(self.menu_add_label("Heure"),self.la_heure_photos)
        self.central_layout.addRow(self.menu_add_label("Écart"),self.le_écart_photos)
        self.central_layout.addRow(self.menu_add_label("Commentaire)"),self.le_commentaire_photos)
        self.central_layout.addRow(self.menu_add_label("Lieu"),self.le_lieu_photos)

        self.central_layout.addRow(self.menu_add_pushbutton("Metre à jour ...",self.PHOTOS_UpdatePhotos_selection, n_cell = 1))
#===========================================================================================================
    def central_import_photos(self):
        self.menu_clear_layout(self.central_layout)       
        self.list_albums(self.central_layout,self.central_import_photos_exec) 
#===========================================================================================================               
    def central_import_photos_exec(self):
        self.menu_clear_layout(self.central_layout)
#----------------------------------------------------------------------------------------------------
        self.la_date_photos = self.menu_add_date(default = self.last_date,ncol = 2)
        self.la_heure_photos = self.menu_add_heure(default = self.last_heure,ncol = 2)
        self.le_écart_photos =self.menu_add_input_text_combobox(["10 minutes", "30 minutes", "1 heure"],n_cell = 1)
        self.le_commentaire_photos =self.menu_add_input_text_simple("No",n_cell = 1)
        self.le_lieu_photos = self.menu_add_input_text_simple("No",n_cell = 1)

        self.central_layout.addRow(self.menu_add_label("Jour"),self.la_date_photos)
        self.central_layout.addRow(self.menu_add_label("Heure"),self.la_heure_photos)
        self.central_layout.addRow(self.menu_add_label("Écart"),self.le_écart_photos)
        self.central_layout.addRow(self.menu_add_label("Commentaire)"),self.le_commentaire_photos)
        self.central_layout.addRow(self.menu_add_label("Lieu"),self.le_lieu_photos)

        self.central_layout.addRow(self.menu_add_pushbutton("Metre à jour ...",self.PHOTOS_Import, n_cell = 1))
#===========================================================================================================   
    def menu_add_layout(self,type):
#--------------------------------------------------------------------------------------------------------------
        if type == "grid": item_layout = QGridLayout()
        else: item_layout = QFormLayout()

        item_layout.setVerticalSpacing(line_spacing)
        item_layout.setHorizontalSpacing(cell_spacing)  
        item = QWidget(self)
        item.setLayout(item_layout)
        return item,item_layout
#============================================================================================================== 
    def menu_clear_layout(self,le_layout):
#--------------------------------------------------------------------------------------------------------------
        while le_layout.count() > 0:
            item = le_layout.itemAt(0)
            widget = item.widget()
            if widget is None:
                le_layout.removeItem(item)
            else:
                #print(widget)
                widget.deleteLater()
                le_layout.removeWidget(widget)
        #print(le_layout.count())
        return
#==============================================================================================================         
    def menu_add_input_text_simple(self,default_input_text,**kwargs):
#--------------------------------------------------------------------------------------------------------------
        if not default_input_text : return
        la_width = None
        n_cell = 2
        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "width" : la_width = valeur
                if clef == "n_cell" : n_cell = valeur
        if not la_width : la_width = n_cell * cell_width
#-------------------------------------------------------------------------------------------------------------- 
        item = QLineEdit(default_input_text)     
        item.setStyleSheet("background-color: "+couleur_femme+";color: black")
        item.setAlignment(Qt.AlignmentFlag.AlignHCenter)
        item.setFont(custom_font) 
        item.setFixedSize(la_width,line_height) 
        return item
#==============================================================================================================
    def menu_add_input_text_completor(self,default_input_list,**kwargs):
#--------------------------------------------------------------------------------------------------------------
        if not default_input_list : return
        la_width = None
        le_alignment = "center"
        n_cell = 2
        popup_action = None
        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "width" : la_width = valeur
                if clef == "n_cell" : n_cell = valeur
                if clef == "popup"  : popup_action = valeur
                if clef == "alignment": le_alignment = valeur
        if not la_width : la_width = n_cell * cell_width

#--------------------------------------------------------------------------------------------------------------    
        le_completor = QCompleter(default_input_list)
        le_completor.setCaseSensitivity(Qt.CaseSensitivity.CaseInsensitive)
        if popup_action : le_completor.popup().clicked.connect(popup_action)
        le_completor.popup().setStyleSheet(
            "QListView {background-color: "+couleur_femme+";color: black;"
                        "selection-background-color: "+couleur_homme+";}"
            )
        item = QLineEdit(self)            
        item.setCompleter(le_completor)
        item.setStyleSheet("background-color: "+couleur_femme+";color: black")
        
        if le_alignment == "left":
            item.setAlignment(Qt.AlignmentFlag.AlignLeft)
        elif le_alignment == "right":
            item.setAlignment(Qt.AlignmentFlag.AlignRight)
        else:
            item.setAlignment(Qt.AlignmentFlag.AlignHCenter)
        item.setFont(custom_font) 
        item.setFixedSize(la_width,line_height) 
        
        return item     
#==============================================================================================================
    def menu_add_input_text_combobox(self,default_input_list,**kwargs):
#--------------------------------------------------------------------------------------------------------------
        if not default_input_list : return
        la_width = None
        n_cell = 2
        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "width" : la_width = valeur
                if clef == "n_cell" : n_cell = valeur
        if not la_width : la_width = n_cell * cell_width
#--------------------------------------------------------------------------------------------------------------
        item = QComboBox()
        item.addItems(default_input_list)      
        item.setStyleSheet("background-color: "+couleur_homme+";color: black")
        item.setFont(custom_font) 
        item.setFixedSize(la_width,line_height) 

        return item
#==============================================================================================================
    def menu_add_infos_row(self,le_texte):
#--------------------------------------------------------------------------------------------------------------
        self.statusbar.showMessage(le_texte)
        self.statusbar.setStyleSheet("background-color: "+couleur_homme+";color: black")
#==============================================================================================================
    def menu_add_label(self,le_label,**kwargs):
#--------------------------------------------------------------------------------------------------------------
        la_width = None
        bcolor = "lightgray"
        color = "black"
        n_cell = 1
        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "width" : la_width = valeur
                if clef == "n_cell" : n_cell = valeur
                if clef == "bcolor": bcolor = valeur
                if clef == "color": color = valeur

        if not la_width : la_width = n_cell * cell_width
#--------------------------------------------------------------------------------------------------------------                   
        item = QLabel(le_label)
        item.setStyleSheet("background-color: "+ bcolor +";color: "+color+";"
                                "qproperty-alignment: AlignLeft;"
                                "border-radius: 10px;"
                                "qproperty-wordWrap: true;"
                                "padding: 6px;"
                                )
        item.setFont(custom_font) 
        #item_label.setAlignment(Qt.AlignmentFlag.AlignHCenter)
        item.setFixedSize(la_width,line_height)

        return item
#==============================================================================================================
    def menu_add_date(self,**kwargs):
        la_width = None
        n_cell = 1
        format_osxphotos = "yyyy-MM-dd"
        format_edit = "dd-MM-yyyy"
        default = False

        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "width" : la_width = valeur
                if clef == "n_cell" : n_cell = valeur
                #if clef == "format":  format = valeur
                if clef == "default" : default = valeur
        if not la_width : la_width = n_cell * cell_width
#--------------------------------------------------------------------------------------------------------------
        item =QDateEdit()
        item.setDisplayFormat(format_edit)
        if default :
            item.setDate(QDate.fromString(default,format_edit))
        else :
            item.setDate(QDate.currentDate())
            
        item.setFont(custom_font)
        item.setFixedSize(la_width, line_height)
        item.setAlignment(Qt.AlignmentFlag.AlignHCenter)

        return item
#==============================================================================================================
    def menu_add_heure(self,**kwargs):
        la_width = None
        n_cell = 1
        default = False

        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "width" : la_width = valeur
                if clef == "n_cell" : n_cell = valeur
                if clef == "default" : default = valeur
        if not la_width : la_width = n_cell * cell_width
#--------------------------------------------------------------------------------------------------------------
        item =QTimeEdit()
        item.setDisplayFormat('HH:mm')
        if default :
            item.setTime(QTime.fromString(default,'HH:mm:ss'))
        else :
            #item.setTime(QTime.currentTime())
            item.setTime(QTime(10,00,00))
            
        item.setFont(custom_font)
        item.setFixedSize(la_width, line_height)
        item.setAlignment(Qt.AlignmentFlag.AlignHCenter)

        return item
#==============================================================================================================
    def menu_add_pushbutton(self,le_texte,cmd,**kwargs):
#--------------------------------------------------------------------------------------------------------------
        la_width = None
        n_cell = 1
        bcolor = "green"
        color = "white"

        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "width" : la_width = valeur
                if clef == "n_cell" : n_cell = valeur
                if clef == "color" : color = valeur
                if clef == "bcolor" : bcolor = valeur
        if not la_width : la_width = n_cell * cell_width
#--------------------------------------------------------------------------------------------------------------
        item = QPushButton(le_texte)
        item.setStyleSheet(
                "QPushButton{background-color : "+bcolor+";color: "+color+"}"
                "QPushButton::pressed{background-color : red ;color: white}"
                        )
        item.setFont(custom_font)
        item.clicked.connect(cmd)
        item.setFixedSize(la_width, line_height)

        return item
#==============================================================================================================
    def menu_add_checkbox(self,le_texte,**kwargs):
#--------------------------------------------------------------------------------------------------------------
        la_width = None
        n_cell = 1
        color_check = couleur_homme
        color = couleur_femme

        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "width" : la_width = valeur
                if clef == "n_cell" : n_cell = valeur
                if clef == "color" : color = valeur
                if clef == "bcolor" : bcolor = valeur
        if not la_width : la_width = n_cell * cell_width
#--------------------------------------------------------------------------------------------------------------
        item = QCheckBox(le_texte)
        item.setStyleSheet(
            "QCheckBox{background-color : "+color+";color:black;}"
            "QCheckBox::checked{background-color : "+color_check+";color:black;}" )

        item.setFont(custom_font)
        item.setFixedSize(la_width, line_height)

        return item
#==============================================================================================================   
    def menu_add_action(self,**kwargs):
#--------------------------------------------------------------------------------------------------------------
        tip = None
        cmd = None
        icon = None
        text = None

        for clef, valeur in kwargs.items():  
            if valeur: 
                if isinstance(clef,str): clef =clef.lower()
                if clef == "tip": tip = valeur
                if clef == "cmd": cmd = valeur
                if clef == "icon": icon = valeur
                if clef == "text": text = valeur
#--------------------------------------------------------------------------------------------------------------

        if icon and text :       item = QAction(QIcon(icon_dir+"/"+icon), text, self)
        elif not icon and text : item = QAction(text, self)
        elif icon and not text : item = QAction(QIcon(icon_dir+"/"+icon), "-",self)
        else: return None

        if cmd : item.triggered.connect(cmd)
        if tip : 
            #item.setToolTip(tip)
            item.setStatusTip(tip)
        item.setCheckable(True)
    
        return item                        
#==============================================================================================================
    def les_sub_folders(self,f_cur):
        temp_list = []

        for le_album in f_cur.album_info:

            if le_album.title[0] != "P": temp_list.append(f' > {le_album.title}')

        if temp_list   :   
            self.list_selection_albums.append(f_cur.title)
            self.list_folder_albums.append([f_cur.title,temp_list])
            self.list_selection_albums = self.list_selection_albums + temp_list

        for sf in f_cur.subfolders:
            self.les_sub_folders(sf)
        return 
#==============================================================================================================           
    def PPTX_reset(self):

        for item in self.personnes_options:
            for iitem in item:
                if iitem:
                    try: iitem.deleteLater() 
                    except : next 
        try:
            self.publication.deleteLater()
            self.mode_couples.deleteLater()
            self.mode_fratrie.deleteLater()
        except : next 
#==============================================================================================================                         
    def PHOTOS_UpdatePhotos_selection(self):
#--------------------------------------------------------------------------------------------------------- 
        format_osxphotos = "yyyy-MM-dd"
        format_edit = "dd-MM-yyyy"

        delta = self.le_écart_photos.currentText()
        if delta == "10 minutes" : d = 600
        elif delta == "30 minutes" : d = 1800


        self.last_date = self.la_date_photos.text()
        str_le_time = self.la_heure_photos.text()
        str_la_date = QDate.toString(QDate.fromString(self.last_date,format_edit),format_osxphotos)
        list_le_date= str_la_date.split("-")

        for idx,photo in enumerate(self.photos_selection):

            if self.le_commentaire_photos.text() != "No" and idx == 0 : photo.description = self.le_commentaire_photos.text()

            if self.le_lieu_photos.text() != "No": 
                ctx = ssl.create_default_context(cafile=certifi.where())
                geopy.geocoders.options.default_ssl_context = ctx
                geolocator = Nominatim(user_agent="my_geocoder")
                location = geolocator.geocode(self.le_lieu_photos.text())
                if location.latitude and location.longitude : photo.location = (location.latitude,location.longitude)

            str_le_time = QTime.toString(QTime.fromString(str_le_time).addSecs(d))
            list_le_time= str_le_time.split(":")

            photo.date = datetime(int(list_le_date[0]),
                                int(list_le_date[1]),
                                int(list_le_date[2]),
                                int(list_le_time[0]),
                                int(list_le_time[1]),
                                int(list_le_time[2])
                                    )


        self.last_heure = QTime.toString(QTime.fromString(str_le_time).addSecs(60))
        
        self.central_updatePhotos_selection()
    #=========================================================================================================
    def PHOTOS_Import(self):
#--------------------------------------------------------------------------------------------------------- 
        location = None
        if self.le_lieu_photos.text() != "No": 
            try: 
                ctx = ssl.create_default_context(cafile=certifi.where())
                geopy.geocoders.options.default_ssl_context = ctx
                geolocator = Nominatim(user_agent="my_geocoder")
                location = geolocator.geocode(self.le_lieu_photos.text())
            except Exception as error: print(error)

        format_osxphotos = "yyyy-MM-dd"
        format_edit = "dd-MM-yyyy"

        delta = self.le_écart_photos.currentText()
        if delta == "10 minutes" : d = 600
        elif delta == "30 minutes" : d = 1800

        self.last_date = self.la_date_photos.text()
        str_le_time = self.la_heure_photos.text()
        str_la_date = QDate.toString(QDate.fromString(self.last_date,format_edit),format_osxphotos)
        list_le_date= str_la_date.split("-")

        import_files = [s for s in os.listdir(import_dir) if os.path.isfile(os.path.join(import_dir, s))]
        import_files.sort(key=lambda s: os.path.getmtime(os.path.join(import_dir, s)))
        import_files_full = [os.path.join(import_dir, s) for s in import_files if s.endswith(".jpg") ]

        photoslib = photoscript.PhotosLibrary()
        import_album = photoslib.album(self.list_selection_albums[-1])
        #imported_photo = photoslib.import_photos([import_files], album=import_album,skip_duplicate_check=True)

        for idx,import_file in enumerate(import_files_full):

            imported_photo_list = photoslib.import_photos([import_file], album=import_album,skip_duplicate_check=True)
            photo = imported_photo_list[0]

        #photoslib.quit()

            if self.le_commentaire_photos.text() != "No" and idx == 0 : photo.description = self.le_commentaire_photos.text()
            if location: photo.location = (location.latitude,location.longitude)

            str_le_time = QTime.toString(QTime.fromString(str_le_time).addSecs(d))
            list_le_time= str_le_time.split(":")

            photo.date = datetime(int(list_le_date[0]),
                                int(list_le_date[1]),
                                int(list_le_date[2]),
                                int(list_le_time[0]),
                                int(list_le_time[1]),
                                int(list_le_time[2])
                                    )

        self.last_heure = QTime.toString(QTime.fromString(str_le_time).addSecs(60))

        self.menu_add_infos_row(f'{len(import_files_full)} photo(s) importée(s)')

        for file in os.listdir(import_dir):
            try: os.remove(import_dir+file)
            except Exception as error: print(error)
        
        self.central_import_photos_exec()
    #=========================================================================================================
    def PDF_generate(self):
#--------------------------------------------------------------------------------------------------------- 
        from fpdf import FPDF

        slide_png_dir = self.PDF_leWIP.currentText()
        s_slide_png_dir = slide_png_dir.split("/")
        le_FileName = s_slide_png_dir[-2]
        #path = "/".join(s_slide_png_dir[0:-2])
        output_pdf_path = les_PDF_RV+le_FileName+"_RV.pdf"

        dirFiles = []
        for f in os.listdir(slide_png_dir):
            if ".png" in f :  dirFiles.append(f)

        if dirFiles :
            dirFiles.sort(key=lambda f: int(''.join(filter(str.isdigit,f))))

            fpdf = FPDF(orientation="landscape", format="A4")
            idx_avant_derniere_image = len(dirFiles) - 1
            for idx,f in enumerate(dirFiles):

                if idx == 0: 
                    n_page = PDF_add_page(fpdf,slide_png_dir+f,0)
                    n_page = PDF_add_page(fpdf,"",n_page)

                else :
                    if idx ==idx_avant_derniere_image and n_page%2 == 0 : n_page = PDF_add_page(fpdf,"",n_page)
                    n_page = PDF_add_page(fpdf,slide_png_dir+f,n_page)

            # Save the PDF document
            fpdf.output(output_pdf_path)
            self.infos.append(f'PDF Recto, Verso > {output_pdf_path} généré')
            self.central_accueil()

        else :
            self.menu_add_infos_row("PDF Recto, Verso > Pas de FileName image")

        return
    #=========================================================================================================                   
    def PPTX_publication_album(self):
#-------------------------------------------------------------------------------------------------------------   
        self.le_document = Presentation(le_model)

        sortby = self.sorttype.currentText()
        isBreakDate = self.is_breakdate.isChecked()
        isBreakCity = self.is_breakcity.isChecked()
        isBreakDescription = self.is_breakdescription.isChecked()

        if not isBreakDate and not isBreakCity and not isBreakDescription : 
            isBreakFree = True
            y_break = 0
        else : 
            isBreakFree = False
            y_break = 9


        isPerson = self.is_person.isChecked()
        isDate = self.is_date.isChecked()
        isCity = self.is_city.isChecked()
        isFilename = self.is_filename.isChecked()
        isDescription = self.is_description.isChecked()
        isTitle = self.is_title.isChecked()
        
        if isBreakDate : isDate = True
        if isBreakCity : isCity = True
        if isBreakDescription : isDescription = True


        if self.marge_haute.currentText() == "Compact" : 
            y_top = 0 + y_break
            y_gap = 0
            n_row_max = 2
            x_gap = 0
        elif self.marge_haute.currentText() == "Dense" : 
            y_top = 0 + y_break
            y_gap = 1
            n_row_max = 3
            x_gap = 1 
        elif self.marge_haute.currentText() == "Classic" : 
            y_top = 2 + y_break
            y_gap = 2
            n_row_max = 2
            x_gap = 2 
        else : exit()

        photo_height = (slide_height - y_top)/n_row_max - y_gap
        #-------------------------------------------------------------------------------------------------------------   
        for idx,le_album_selection in enumerate(self.list_selection_albums):
        #-------------------------------------------------------------------------------------------------------------
            photos = self.photosdb.photos(albums= [le_album_selection]) 
            register_heif_opener()

            if photos:  
                buffer_photos = []
                les_erreurs = []
                for idx,photo in enumerate(photos):
                #----------------------------------------------------------------------------------------------------
                    if photo.isphoto:                         

                        # sort key
                        if sortby  == "FileName":  la_clef = photo.original_filename 
                        elif sortby  == "Title": 
                            if photo.title : la_clef = photo.title
                            else: la_clef = photo.date.strftime('%Y%m%d')
                        else : la_clef = photo.date.strftime('%Y%m%d%H%M%S')

                        #initiate my_photo object
                        p           = my_photo()
                        p.id        = photo.uuid
                        p.filename  = photo.original_filename
                        p.date      = photo.date.strftime('%d-%m-%Y')
                        p.datelong  = traduction_date(photo.date.strftime('%A %d %B %Y'))
                        p.favorite  = photo.favorite

                        # Persons
                        temp_persons = []
                        for item in photo.persons:
                            if item != "_UNKNOWN_" : temp_persons.append(item)
                        if temp_persons : p.person = ", ".join(temp_persons)
                        #Places 
                        if photo.place:
                            if photo.place.names.area_of_interest : p.poi       =  photo.place.names.area_of_interest[0]
                            if photo.place.names.street_address :   p.street    =  photo.place.names.street_address[0]
                            if photo.place.names.city :             p.city      =  photo.place.names.city[0]
                            if photo.place.names.country :          p.country   =  photo.place.names.country[0] 

                        if photo.description:                       p.description = photo.description
                        if photo.title :                            p.title     = photo.title

                        #path
                        hasRgb_im = False
                        for filephoto in photo.path_derivatives:
                            rgb_im = cv2.imread(filephoto)
                            try:  
                                len(rgb_im)
                                hasRgb_im = True
                                break
                            except : next

                        if hasRgb_im : 
                            image_height, image_width = rgb_im.shape[:2]
                            p.width = image_width
                            p.height = image_height
                            photo_filename = f'{export_path}/{idx}.jpg'
                            cv2.imwrite(photo_filename, rgb_im)
                            p.path = photo_filename
                            buffer_photos.append([la_clef,p])
                        else:
                            print(f"{photo.original_filename} path_derivatives problem")
                    else: les_erreurs.append(f"{photo.original_filename} not an image")
                        
                for item in les_erreurs:
                    print(item)

                # Add photos to document
                #-----------------------------------------------------------------------------------------------------------------------------------------
                if buffer_photos:
                    buffer_photos = sorted(buffer_photos, key=lambda col: (col[0]) )

                    #page de garde 
                    #-------------------------------------------------------------------
                    page1,boxes = PPTX_add_page(self.le_document,slide_layout_garde)

                    temp_texte = le_album_selection.replace("LIVRE_","")
                    t_zorg = temp_texte.split("@")
                    if len(t_zorg) > 1: temp_texte_place= t_zorg[-1]
                    else: temp_texte_place = ""
                    temp_texte_date = ""
                    temp_texte_date = t_zorg[0].split(" ")[0]
                    temp_texte = " ".join(t_zorg[0].split(" ")[1:])

                    p = PPTX_add_paragraph(boxes[0])
                    if temp_texte : PPTX_add_run(p,temp_texte)
                    if temp_texte_date : 
                        p = PPTX_add_paragraph(boxes[0])
                        PPTX_add_run(p,temp_texte_date)
                    if temp_texte_place:
                        p = PPTX_add_paragraph(boxes[0])
                        PPTX_add_run(p,temp_texte_place)

                    # initialisation des dimensions
                    #-------------------------------------------------------------------
                    n_row = 1
                    x_bandeau = x_gap
                        
                    x = slide_width - x_gap
                    y = y_top

                    la_page = False
                    isPreviousFavorite = False
                    isNext_row = False
                    paragraph_title_text = ""
                    date_previous = "NEW_PAGE"
                    city_previous = "NEW_PAGE"
                    description_previous = "NEW_PAGE"
                    isBreak = True
                    pic_previous = None

                    # loop on sorted photo buffer
                    #-------------------------------------------------------------------
                    isNewPage = False
                    for idx,buffer_photo in enumerate(buffer_photos):
                    #-------------------------------------------------------------------
                        p = buffer_photo[1] #my_photo object
                        y_diff = 0
                        
                        # resize image
                        percent_crop_horizontal = 0.0
                        percent_crop_vertical = 0.0

                        r = p.width / p.height
                        img_height = photo_height
                        img_width = photo_height * r 

                        # is it a New page ?
                        if  (   (not la_page) or
                                (isBreakDate and p.date != date_previous) or 
                                (isBreakCity and p.city != city_previous) or 
                                (isBreakDescription and  p.description != description_previous)
                            ) : isNewPage = True  

                        # break info title
                        le_temp = []
                        paragraph_title_text = ""
                        if p.description and isBreak and isBreakDescription: paragraph_title_text= p.description
                        if p.date and isBreak and isBreakDate : le_temp.append(p.datelong)
                        if p.city and isBreak and isBreakCity : le_temp.append(p.city)
                        if le_temp: paragraph_title_text = ", ".join(le_temp)    
                    
                        # photo favorite 
                        if p.favorite: 
                            iw = p.width
                            ih = p.height
                            if iw > ih : 
                                img_height = slide_height - y_top -y_gap
                                img_width = img_height * r
                                x_diff =  img_width - slide_width + 2 * x_gap
                                if x_diff > 0: percent_crop_horizontal = x_diff/img_width                                 
                            else:
                                img_width = slide_width  - 2 * x_gap
                                img_height = img_width / r 
                                y_diff =  img_height - slide_height + y_top + y_gap
                                if y_diff > 0: percent_crop_vertical = -y_diff/img_height   

                            isNewPage = True  
                            isPreviousFavorite = True 
                        elif  isPreviousFavorite:
                            isPreviousFavorite = False
                            isNewPage = True
                        

                        if isNewPage:
                        #---New page-------------------------------------------------------------------------------                            
                            if isBreakFree: le_album_layout = slide_layout_album_breakfree
                            else: le_album_layout = slide_layout_album_break
                            la_page, boxes = PPTX_add_page(self.le_document,le_album_layout)
                            if not isBreakFree : PPTX_add_run (PPTX_add_paragraph(boxes[0]),paragraph_title_text)

                            isNewPage = False
                            isNext_row = False
                            isBreak = True
                            x = x_gap
                            y = y_top
                            n_row = 1
                        #--- End New page--------------------------------------------------------------------------

                        #Current page  
                        else:

                            x_diff = int( slide_width - x  - img_width - x_bandeau)

                            if x_diff  < 0 : 

                                if slide_width - x - x_gap - x_bandeau > 30 :
                                    percent_crop_horizontal = -1 * (x_diff / img_width)
                                    isNext_row = True
                                else:
                                    isNext_row = False
                                    PPTX_expand_picture(pic_previous,x_bandeau)
                                    n_row = n_row + 1
                                    
                                    if n_row > n_row_max :
                                    #---New page-------------------------------------------------------------------------------                            
                                        if isBreakFree: le_album_layout = slide_layout_album_breakfree
                                        else: le_album_layout = slide_layout_album_break
                                        la_page, boxes = PPTX_add_page(self.le_document,le_album_layout)
                                        if not isBreakFree : PPTX_add_run (PPTX_add_paragraph(boxes[0]),paragraph_title_text)

                                        isNewPage = False
                                        isNext_row = False
                                        isBreak = True
                                        x = x_gap
                                        y = y_top
                                        n_row = 1
                                    #--- End New page--------------------------------------------------------------------------

                                    else :
                                        x = x_gap
                                        y = y + img_height + y_gap   
    
                        # Add image   
                        xp = x

                        if img_width > 0 and img_height > 0 :
                            #crop image
                            pic = la_page.shapes.add_picture(p.path,Mm( xp ),Mm(y - y_diff),Mm(img_width),Mm(img_height))
                            pic.crop_left = percent_crop_horizontal
                            pic.crop_right = -1 * percent_crop_horizontal
                            pic.crop_top = percent_crop_vertical
                            pic.crop_bottom = -1 * percent_crop_vertical
                            pic_previous = pic

                        else:
                            print("img_width=",img_width, "img_height=",img_height)
                            exit()

                        # Info box management -------------------------------------------------------------
                        iw = img_width * (1 - percent_crop_horizontal)
                        if iw < 10 : info_box_width = 20
                        else :info_box_width = iw 

                        mt = mb = ml = mr = 1

                        font_size = 16
                        
                        info_box_texts = []
                        info_box_titre_temp = []

                        if p.date:
                            if not isBreakDate and isDate and p.date != date_previous :
                                    info_box_titre_temp.append(f'{p.date}')

                        if p.city:
                            if not isBreakCity and isCity and p.city != city_previous : 
                                    info_box_titre_temp.append(f'{p.city}')

                        if info_box_titre_temp: 
                            info_box_texts.append(", ".join(info_box_titre_temp))

                        if p.description:
                            if not isBreakDescription and isDescription and p.description != description_previous:
                                    info_box_texts.append(p.description)

                        if isTitle and p.title :info_box_texts.append(p.title)
                        if isFilename and p.filename : info_box_texts.append(p.filename)
                        if isPerson and p.person : info_box_texts.append(p.person)

                        if info_box_texts:
                            info_box_height = box_height(info_box_texts,info_box_width,"normal",font_size,mt,mb,ml,mr)

                            object_textbox = PPTX_add_box(la_page,xp,y,info_box_width , info_box_height,
                                                            "transparent","autosize",
                                                            margin_bottom = mb,margin_top = mt,margin_left= ml,margin_right = mr,
                                                            bcolor=Darkblue)
                            for idx,item in enumerate(info_box_texts):
                                if item : PPTX_add_run(PPTX_add_paragraph(object_textbox),item,size = font_size,color = White) 


                        date_previous = p.date
                        city_previous = p.city
                        description_previous = p.description
                        # End Info box management ----------------------------------------------------------
                        
                        # prepare new photo
                        if isNext_row :
                                n_row = n_row + 1
                                
                                if n_row > 2 :
                                    isNewPage = True
                                else :
                                    x = x_gap
                                    y = y + photo_height + y_gap  
                                    isNext_row = False
                        else:
                            x = x + img_width + x_gap 

                    #-------------------------------------------------------------------------------------------------------------
                    try:
                        for filename in glob.iglob(f"{export_path}/*.*", recursive=False):
                            os.remove(filename)
                    except Exception as error:
                        print("PPTX_publication_album",error)
                    #-------------------------------------------------------------------------------------------------------------
                    self.les_infos = f"+ Album {le_album_selection} {sortby}"
                    self.menu_add_infos_row("Albums > "+self.les_infos)

                    #ZOB

                    #try:
                    #    le_FileName = self.OUTFileName.text()
                    #    os.remove(le_FileName)
                    #except: next

                    self.PPTX_save()




#==============================================================================================================   
#==============================================================================================================                           
#==============================================================================================================   
if __name__ == "__main__":
    import sys
    app = QApplication(sys.argv)
    MainWindow = LA_WINDOW()
    MainWindow.show()
    sys.exit(app.exec())
#==============================================================================================================   
#==============================================================================================================                           
#==============================================================================================================   
